"""模組 3：AI 物件行銷工廠 - 10 個實戰案例 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 品牌色
NAVY = RGBColor(0x0A, 0x16, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CARD_BG = RGBColor(0x0F, 0x1D, 0x32)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 0.8, 11, 1, "模組 3", 28, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "📢 AI 物件行銷工廠", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 1, "10 個實戰案例 ─ 一間物件，14 件跨平台內容", 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, "AI 超級房仲實戰班 ｜ 講師：阿峰老師（黃敬峰）", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "www.autolab.cloud", 14, CYAN, False, PP_ALIGN.CENTER)

# ========== 案例資料 ==========
cases = [
    {
        "num": "01",
        "title": "591 刊登標題（吸睛版）",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：物件行銷一鍵產出包的第一站",
        "scenario": "你剛接到一間板橋 2 房委售案，要上 591 刊登。標題是買方搜尋時第一眼看到的東西，寫得好不好直接影響點擊率。",
        "task": "上傳物件資料，請 AI 產出 3 版 591 標題 + 關鍵字標籤",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "整理好物件基本資料（坪數、格局、屋齡、特色、開價等）",
        "prompt": '我是房仲，以下是一間委售物件的資料。\n\n請幫我產出 3 版 591 刊登標題：\n- 版本 A：吸睛版（用數字或問句吸引點擊）\n- 版本 B：專業版（強調地段、坪效、增值潛力）\n- 版本 C：溫馨版（強調生活感、適合小家庭）\n\n每個標題控制在 30 字以內。\n\n另外請產出 10 個搜尋關鍵字標籤（買方可能會搜尋的詞），例如：板橋、新板特區、2房、首購、近捷運等。',
    },
    {
        "num": "02",
        "title": "591 物件描述（完整版）",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "591 標題寫好了，接下來要寫 500-800 字的詳細描述。好的描述能讓買方看完就想約看房，而不是滑過去。",
        "task": "上傳物件資料，請 AI 產出 591 完整物件描述",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "同案例 01 的物件資料，讓 AI 從同一份資料延伸產出描述",
        "prompt": '我是房仲，以下是委售物件的資料。\n\n請幫我寫一篇 591 物件描述（500-800 字），結構如下：\n1. 開頭吸引（用一句話讓買方想繼續看）\n2. 五大亮點（格局、採光、交通、生活機能、社區管理）\n3. 生活機能描述（周邊學校、市場、公園、商圈）\n4. 交通便利性（捷運站、公車、開車動線）\n5. 數據佐證（實價登錄、區域漲幅、租金投報率）\n6. CTA 結尾（製造緊迫感，引導聯繫看房）\n\n語氣要專業但親切，避免過度誇大。\n重要：不要寫「絕對增值」「保證賺錢」等違規用語。',
    },
    {
        "num": "03",
        "title": "LINE 圖卡文案",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "要做一張 LINE 圖卡傳給客戶群組。圖卡要一眼就看到重點，文字不能太多但要夠吸引人。",
        "task": "上傳物件資料，請 AI 產出圖卡的主標 + 副標 + 條列亮點 + CTA",
        "tool": "ChatGPT + Canva",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "產出文案後，直接貼到 Canva 的房仲圖卡模板中",
        "prompt": '我是房仲，以下是物件資料。\n\n請幫我設計一張 LINE 圖卡的文案，要求：\n1. 主標題（15 字以內，要有記憶點）\n2. 副標題（25 字以內，補充關鍵賣點）\n3. 3 個條列亮點（每點 10 字以內，用 emoji 開頭）\n4. CTA 按鈕文案（8 字以內，例如「立即預約看房」）\n5. 底部資訊（房仲姓名 + 電話）\n\n文案風格：簡潔有力，一眼就能抓到重點。\n產出 2 個版本讓我選：版本 A 走專業風、版本 B 走溫馨風。\n\n這些文案要能直接貼到 Canva 模板中使用。',
    },
    {
        "num": "04",
        "title": "物件推薦通知訊息",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：新物件上架自動推播給匹配客戶",
        "scenario": "新委售物件上架了，你手上有一批客戶名單。不同客戶的需求不同（首購、換屋、投資），不能用同一則訊息群發。",
        "task": "上傳物件資料和客戶名單，請 AI 針對不同客戶類型產出推薦訊息",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔 + 📄 CSV 檔案",
        "upload_file": "物件資料.txt + 客戶名單.csv",
        "upload_note": "CSV 欄位建議：姓名、需求類型（首購/換屋/投資）、預算、偏好區域",
        "prompt": '我是房仲，剛接到一間新委售物件，以下是物件資料。\n另外附上我的客戶名單 CSV。\n\n請幫我：\n1. 從客戶名單中篩選出可能對這間物件有興趣的客戶\n2. 針對 3 種客戶類型，各寫一則推薦 LINE 訊息：\n   - 首購族：強調總價親民、首購優勢、生活機能\n   - 換屋族：強調格局升級、社區品質、學區\n   - 投資客：強調租金投報率、增值潛力、出租容易度\n\n每則訊息 100 字以內，語氣自然像一對一私訊。\n不要用群發感的語氣（例如「各位朋友大家好」）。',
    },
    {
        "num": "05",
        "title": "FB 故事行銷貼文",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "要在 FB 粉專發一篇故事行銷貼文。不是硬梆梆的物件資訊，而是要用「理想生活」的角度讓人產生嚮往。",
        "task": "上傳物件資料和照片，請 AI 把物件變成有故事性的 FB 貼文",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔 + 🖼️ 物件照片",
        "upload_file": "物件資料.txt + 物件照片",
        "upload_note": "上傳 2-3 張物件最好看的照片，讓 AI 也能參考空間感來寫文案",
        "prompt": '我是房仲，以下是物件資料和照片。\n\n請幫我寫一篇 FB 粉專貼文（300-500 字），要求：\n1. 以「理想生活」切入，不要開頭就寫「稀有釋出」「難得好房」\n2. 用一個生活場景開頭（例如：「週六早上，推開窗，陽光灑進客廳...」）\n3. 把物件亮點融入故事中，而不是條列式\n4. 結尾製造行動力（但不要太業務感）\n5. 附上 10 個 hashtag（含區域、房型、生活風格標籤）\n\n語氣要求：\n- 像在跟朋友分享一個好地方\n- 不要用「限時」「僅此一間」等高壓用語\n- 要讓人覺得「這就是我想要的生活」',
    },
    {
        "num": "06",
        "title": "IG 限時動態腳本（5 則）",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "要在 IG 限動連發 5 則介紹物件。限動的特性是快速、互動、一則一重點，要讓人看到第 5 則就想私訊。",
        "task": "上傳物件資料和照片，請 AI 設計 5 則限動的文案和互動元素",
        "tool": "ChatGPT + Canva",
        "upload_format": "📄 文字檔 + 🖼️ 物件照片 5 張",
        "upload_file": "物件資料.txt + 物件照片 5 張",
        "upload_note": "建議照片順序：外觀、客廳、主臥、廚房/衛浴、周邊環境",
        "prompt": '我是房仲，以下是物件資料和 5 張照片。\n\n請幫我設計 5 則 IG 限時動態的腳本：\n\n第 1 則（封面吸引）：\n- 大字標題（8 字以內）\n- 互動：投票貼紙「你覺得這間值多少？」\n\n第 2 則（格局介紹）：\n- 重點文字標註在照片上的位置建議\n- 互動：問答貼紙「猜猜幾坪？」\n\n第 3 則（生活機能）：\n- 周邊亮點 3 個（學校/捷運/商圈）\n- 互動：倒數貼紙「開放預約看房」\n\n第 4 則（價格透明）：\n- 總價 + 每坪單價 + 自備款試算\n- 互動：滑桿貼紙「這個價格你覺得？」\n\n第 5 則（CTA）：\n- 「想看房？私訊我」\n- 互動：連結貼紙（導到 LINE 或預約頁）\n\n每則含：文字內容 + 建議貼紙 + 互動設計。',
    },
    {
        "num": "07",
        "title": "IG 輪播貼文文案",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "要做一組 6 張的 IG 輪播貼文。輪播的重點是「滑到最後一張就想行動」，每張都要有看下去的動力。",
        "task": "上傳物件資料，請 AI 設計 6 張輪播的標題和內容",
        "tool": "ChatGPT + Canva",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "輪播文案產出後，搭配 Canva 輪播模板使用",
        "prompt": '我是房仲，以下是物件資料。\n\n請幫我設計 6 張 IG 輪播貼文：\n\n第 1 張（封面）：\n- 標題要有鉤子（例如問句或數字）\n- 副標一句話說明物件亮點\n\n第 2 張（亮點一）：\n- 格局 / 採光 / 空間感\n- 標題 + 2-3 行說明\n\n第 3 張（亮點二）：\n- 交通 / 生活機能\n- 標題 + 2-3 行說明\n\n第 4 張（亮點三）：\n- 社區 / 管理 / 附加價值\n- 標題 + 2-3 行說明\n\n第 5 張（數據頁）：\n- 總價 / 坪數 / 每坪單價 / 屋齡\n- 同區比較一句話（例如「比同社區便宜 X%」）\n\n第 6 張（CTA）：\n- 「想看房？存起來 + 私訊我」\n- 附上聯繫方式\n\n每張包含：標題（15字內）+ 內文（50字內）。\n另外寫一段 IG 貼文描述（caption，200字內 + hashtag）。',
    },
    {
        "num": "08",
        "title": "節日借勢行銷貼文",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：全年節日行事曆自動觸發",
        "scenario": "農曆新年快到了，你想結合節日氛圍來推物件。節日行銷的關鍵是「讓物件搭上節日的情感」，不是硬塞節日元素。",
        "task": "上傳物件資料，請 AI 寫出結合節日氛圍的行銷貼文",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "可替換不同節日（新年/端午/中秋/聖誕），AI 會自動調整氛圍",
        "prompt": '我是房仲，以下是物件資料。農曆新年快到了。\n\n請幫我寫結合新年主題的行銷貼文：\n\n1. FB 貼文（300-400 字）：\n   - 用「新年新家」「開運好宅」的概念切入\n   - 把物件亮點跟新年意象連結（例如：採光好 = 迎接新年第一道陽光）\n   - 結尾用新年祝福收尾，自然帶到看房邀約\n   - 10 個 hashtag（含節日 + 房產標籤）\n\n2. IG 貼文（200 字 + hashtag）：\n   - 比 FB 更簡短活潑\n   - 開頭用新年相關的 emoji 吸引目光\n   - 亮點用條列式（3 點）\n\n3. LINE 圖卡文案：\n   - 主標：8 字以內（含新年元素）\n   - 副標：15 字以內\n   - CTA：「新年新家，立即預約」\n\n語氣溫暖喜氣，但不要太浮誇。\n物件亮點要自然融入節日主題，不是硬湊。',
    },
    {
        "num": "09",
        "title": "Google 商家成交實績貼文",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：成交後自動產出 Google 商家貼文，提升 SEO",
        "scenario": "案子成交了！要在 Google 商家檔案發一篇成交實績貼文。這不只是分享喜悅，更是提升在地 SEO 排名的重要動作。",
        "task": "上傳成交資訊，請 AI 寫出專業的成交實績貼文",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔",
        "upload_file": "成交資訊（不含客戶隱私）.txt",
        "upload_note": "只提供區域、房型、價格帶，不要包含客戶姓名和詳細地址",
        "prompt": '我是房仲，剛成交一間物件，以下是成交資訊（已去除客戶隱私）。\n\n請幫我寫一篇 Google 商家成交實績貼文，要求：\n1. 開頭：恭喜成交（喜氣但專業）\n2. 服務過程簡述（2-3 句，展現專業度）\n3. 強調該區域的專業度（讓搜尋這個區域的人看到）\n4. 邀請客戶留下 Google 評論\n5. 結尾 CTA：「如果你也在找 [區域] 的房子...」\n\nSEO 關鍵字要求：\n- 自然置入區域名稱至少 3 次\n- 包含「房仲推薦」「成交實績」等關鍵字\n- 200-300 字，不要太長\n\n另外請產出：\n- 一則可以轉發到 FB 的版本\n- 一則精簡版（適合 LINE 分享，80 字以內）',
    },
    {
        "num": "10",
        "title": "YouTube Shorts 60 秒腳本",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：物件行銷一鍵產出包",
        "scenario": "要拍一支 60 秒的 YouTube Shorts 介紹物件。短影片的關鍵是「前 3 秒留住人」，然後每 10 秒一個新資訊，最後 CTA。",
        "task": "上傳物件資料，請 AI 寫出完整的 60 秒短影片腳本",
        "tool": "ChatGPT",
        "upload_format": "📄 文字檔",
        "upload_file": "物件資料.txt",
        "upload_note": "腳本產出後，可以直接照著拍。建議搭配手機穩定器",
        "prompt": '我是房仲，以下是物件資料。\n\n請幫我寫一支 60 秒 YouTube Shorts 腳本，結構如下：\n\n0-3 秒（鉤子）：\n- 一句話抓住注意力（例如：「板橋這間 2 房，我帶看 3 組都搶著要出價」）\n- 鏡頭：面對鏡頭說話\n\n3-10 秒（問題切入）：\n- 點出買方的痛點（例如：「在板橋找 2 房，預算 1,500 以內，幾乎不可能？」）\n- 鏡頭：邊走邊說，走進物件大門\n\n10-35 秒（物件亮點巡覽）：\n- 3 個主要亮點，每個 8 秒\n- 每個亮點含：旁白 + 字卡文字 + 鏡頭建議\n\n35-50 秒（數據佐證）：\n- 價格、坪數、同區比較\n- 字卡顯示關鍵數字\n\n50-60 秒（CTA）：\n- 「想看房？下面留言或私訊我」\n- 字卡：聯繫方式\n- 鏡頭：面對鏡頭，親切收尾\n\n每個段落都要標註：秒數 / 旁白 / 字卡 / 鏡頭指示。',
    },
]

# ========== 生成案例投影片 ==========
for case in cases:
    # 情境頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 頂部標籤（含 level）
    add_text_box(slide, 0.5, 0.3, 5, 0.5, f"案例 {case['num']}  {case['level']}", 16, CYAN, True)
    add_text_box(slide, 0.5, 0.7, 12, 0.8, case["title"], 36, WHITE, True)

    # 左半邊：情境 + 任務
    tf = add_text_box(slide, 0.5, 1.8, 6, 0.5, "📋 情境", 20, ORANGE, True)
    add_paragraph(tf, case["scenario"], 17, WHITE, False, 8)
    add_paragraph(tf, "", 10, WHITE, False, 12)
    add_paragraph(tf, "🎯 任務目標", 20, ORANGE, True, 8)
    add_paragraph(tf, case["task"], 17, WHITE, False, 8)

    # 右半邊：工具 + 上傳資料
    tf2 = add_text_box(slide, 7, 1.8, 6, 0.5, "🔧 使用工具", 20, ORANGE, True)
    add_paragraph(tf2, case["tool"], 17, CYAN, True, 8)
    add_paragraph(tf2, "", 10, WHITE, False, 12)
    add_paragraph(tf2, "📎 要準備的資料", 20, ORANGE, True, 8)
    add_paragraph(tf2, f"格式：{case['upload_format']}", 15, GRAY, False, 8)
    add_paragraph(tf2, f"檔案：{case['upload_file']}", 15, WHITE, False, 4)
    add_paragraph(tf2, f"💡 {case['upload_note']}", 13, GRAY, False, 8)
    if case.get('system_note'):
        add_paragraph(tf2, case['system_note'], 13, CYAN, False, 10)

    # 操作頁（Prompt）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY)

    add_text_box(slide2, 0.5, 0.3, 6, 0.5, f"案例 {case['num']} ─ 操作指令", 16, CYAN, True)
    add_text_box(slide2, 0.5, 0.7, 12, 0.8, case["title"], 28, WHITE, True)

    # Prompt 區塊
    add_rounded_rect(slide2, 0.5, 1.6, 12.3, 5.2, CARD_BG)
    tf3 = add_text_box(slide2, 0.8, 1.7, 11.8, 0.4, "📝 複製以下指令，貼到 AI 工具中：", 16, ORANGE, True)

    prompt_lines = case["prompt"].split("\n")
    for line in prompt_lines:
        add_paragraph(tf3, line, 14, WHITE, False, 2)

    add_text_box(slide2, 0.5, 6.9, 12, 0.4, "💡 Prompt 可依實際物件修改關鍵字，不需要逐字照抄", 13, GRAY)

# ========== SOP 流程圖頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "物件行銷工廠 SOP", 32, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.2, 12, 0.5, "一間物件，14 件跨平台內容素材", 18, GRAY, False, PP_ALIGN.CENTER)

sop_lines = [
    "接到新委售 → 整理物件資料.txt",
    "    ↓",
    "案例 01 → 591 標題 × 3 版",
    "案例 02 → 591 完整描述",
    "    ↓",
    "案例 03 → LINE 圖卡文案 → Canva 出圖",
    "案例 04 → 物件推薦通知 → 推播給匹配客戶",
    "    ↓",
    "案例 05 → FB 故事貼文",
    "案例 06 → IG 限動 5 則",
    "案例 07 → IG 輪播 6 張",
    "案例 10 → YT Shorts 腳本",
    "    ↓",
    "成交後 → 案例 09 Google 商家實績貼文",
    "節日時 → 案例 08 借勢行銷貼文",
    "",
    "一間物件 = 14+ 件跨平台內容素材",
]
tf_sop = add_text_box(slide, 1, 2.0, 11, 4.5, "", 15, WHITE)
for line in sop_lines:
    c = CYAN if "↓" in line else (ORANGE if "一間物件 =" in line else WHITE)
    add_paragraph(tf_sop, line, 15, c, "一間物件 =" in line, 3)

add_text_box(slide, 0.5, 6.5, 12, 0.5, "💡 不是內容做不出來，是你沒有系統化。AI 幫你批量產出，你負責選最好的版本", 16, ORANGE, True, PP_ALIGN.CENTER)

# ========== 結尾頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 2, 11, 1, "模組 3 完成！", 44, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "你已經學會用 AI 做物件行銷工廠", 28, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1.5, "記住：一間物件，14 件內容\nAI 產出素材，你負責選擇最對的那版\n\n💡 進階挑戰：把 01→02→03→05→06→07→10 串成「物件行銷一鍵產出包」", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6, 11, 0.5, "講師：阿峰老師 ｜ www.autolab.cloud", 16, CYAN, False, PP_ALIGN.CENTER)

output = "/Users/huangjingfeng/Desktop/專案/02-培訓(Training)/02-公開班(Public-Class)/03-ai-realtor-course(房仲AI課程)/cases/模組3-AI物件行銷工廠-10案例.pptx"
prs.save(output)
print(f"完成：{output}")
print(f"共 {len(prs.slides)} 張投影片")
