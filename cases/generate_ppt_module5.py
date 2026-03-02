"""模組 5：AI 公域獲客飛輪 - 10 個實戰案例 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 品牌色
NAVY = RGBColor(0x0A, 0x16, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CARD_BG = RGBColor(0x0F, 0x1D, 0x32)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 0.8, 11, 1, "模組 5", 28, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "🚀 AI 公域獲客飛輪", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 1, "10 個實戰案例 ─ 社群內容批量生產，永不斷更", 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, "AI 超級房仲實戰班 ｜ 講師：阿峰老師（黃敬峰）", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "www.autolab.cloud", 14, CYAN, False, PP_ALIGN.CENTER)

# ========== 案例資料 ==========
cases = [
    {
        "num": "01",
        "title": "LINE 早安房市快訊（批量20則）",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：社群內容批量生產系統 Batch 1",
        "scenario": "你想每天早上在LINE群組發一則房市小知識，但沒時間每天寫。如果能一次產好整個月的內容，每天早上複製貼上就好了。",
        "task": "一次產出20則早安快訊，撐4週",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "不需要準備資料，直接跟 AI 對話即可",
        "prompt": '我是房仲，想每天早上在 LINE 群組發一則早安房市快訊。\n\n請幫我一次產出 20 則早安快訊，要求：\n1. 每則 50-80 字\n2. 開頭有早安問候語（不重複）\n3. 每則包含一個實用知識點\n4. 涵蓋以下主題（每類至少 3-4 則）：\n   - 房市政策新知\n   - 房貸利率小知識\n   - 裝修實用建議\n   - 居家風水小提醒\n   - 買賣法規提醒\n5. 結尾加一句互動引導（例如「你覺得呢？」「有遇過嗎？」）\n6. 適量 emoji，不要太花俏\n\n請用表格呈現：編號 / 主題分類 / 快訊內容 / 建議發送日（週一到週五）',
    },
    {
        "num": "02",
        "title": "FB 帶看心得分享文",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "你今天帶了 3 組客戶看房，過程中有一些有趣的觀察和市場洞察，想寫成 FB 貼文分享，建立專業形象。",
        "task": "產出帶看分享 FB 貼文",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "把今日帶看經歷直接打字描述給 AI 即可",
        "prompt": '我是房仲，今天帶了 3 組客戶看房，想寫一篇 FB 帶看心得分享文。\n\n今日帶看概況：\n- 客戶 A：新婚夫妻看板橋 2 房，很喜歡但卡預算\n- 客戶 B：投資客看中和套房，問了很多報酬率的問題\n- 客戶 C：退休夫妻看電梯大樓，最在意無障礙和醫院距離\n\n請幫我寫一篇 FB 貼文，要求：\n1. 不暴露客戶隱私（不提真名和具體地址）\n2. 從帶看經歷中提煉 1-2 個專業觀察\n3. 加入市場洞察（例如：最近什麼類型的物件最搶手）\n4. 結尾放一個互動問句（引導留言）\n5. 附上 5-8 個相關 hashtag\n6. 語氣：專業但親切，像在跟朋友聊天\n7. 字數：300-500 字\n\n產出 2 個版本：\n- 版本 A：故事型（用一個客戶的故事帶出觀點）\n- 版本 B：知識型（用帶看經驗分享購屋技巧）',
    },
    {
        "num": "03",
        "title": "IG 個人品牌經營（自我介紹）",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "你要重新整理 IG 帳號，做一系列個人品牌內容，讓潛在客戶透過 IG 認識你、信任你。",
        "task": "設計 IG 個人品牌系列貼文",
        "tool": "ChatGPT + Canva",
        "upload_format": "📷 個人照片 + 📄 經歷資料",
        "upload_file": "個人照片+經歷簡述",
        "upload_note": "上傳個人照片和經歷讓 AI 幫你規劃品牌內容",
        "prompt": '我是房仲，想經營 IG 個人品牌。\n\n我的基本資料：\n- 入行 3 年，成交 50 戶\n- 專精板橋、中和區域\n- 擅長首購族服務\n- 個人特色：耐心、數據分析派\n\n請幫我設計 5 組 IG 貼文內容：\n\n1. 自我介紹（讓人 3 秒記住你）\n   - 標題、內文、hashtag\n\n2. 服務理念（為什麼選你？）\n   - 標題、內文、hashtag\n\n3. 成交故事（一個感人/有趣的成交案例）\n   - 標題、內文、hashtag\n\n4. 購屋知識（提供價值，建立專業感）\n   - 標題、內文、hashtag\n\n5. 客戶見證（社交證明）\n   - 標題、內文、hashtag\n\n每組貼文包含：\n- IG 圖片上的標題文字（10 字以內，Canva 用）\n- 內文（150-200 字）\n- hashtag 策略（大中小標籤各 3 個）\n- 最佳發布時間建議',
    },
    {
        "num": "04",
        "title": "Threads 房產冷知識（批量10則）",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：社群內容批量生產系統 Batch 2",
        "scenario": "你開了 Threads 帳號，想發短小精悍的房產知識，靠冷知識引起好奇心，吸引追蹤者。",
        "task": "批量產出 10 則 Threads 貼文",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "不需要準備資料，直接跟 AI 對話即可",
        "prompt": '我是房仲，想在 Threads 上發短篇房產冷知識來吸粉。\n\n請幫我產出 10 則 Threads 貼文，要求：\n1. 每則 100-150 字\n2. 開頭用一個反直覺的事實或數據抓住注意力\n   （例如：「你知道嗎？台灣有 30% 的房子其實是空的」）\n3. 中間簡短解釋\n4. 結尾引導討論（讓人想回覆）\n5. emoji 適量（每則 2-3 個）\n6. 涵蓋以下主題：\n   - 買房冷知識 × 3\n   - 房貸秘密 × 2\n   - 裝修地雷 × 2\n   - 法規陷阱 × 2\n   - 風水趣聞 × 1\n\n請用表格呈現：編號 / 主題 / 貼文內容 / 預估互動度（高/中/低）',
    },
    {
        "num": "05",
        "title": "電子報月刊設計",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：月報多版本產出的電子報版",
        "scenario": "你想每月發一份電子報給過去的客戶和潛在客戶，維繫關係的同時展現專業。",
        "task": "設計電子報的內容架構和文案",
        "tool": "ChatGPT",
        "upload_format": "📄 當月行情數據",
        "upload_file": "當月行情數據（實價登錄/新聞摘要）",
        "upload_note": "上傳當月的行情數據，讓 AI 結合數據寫出更有說服力的內容",
        "prompt": '我是房仲，想每月發一份電子報給客戶名單（約 200 人）。\n\n請幫我設計這期電子報的完整內容：\n\n1. 開頭：編輯觀點（150 字）\n   - 用一個故事或觀察開場\n   - 帶出本月的主題\n\n2. 市場數據速報（3 個重點）\n   - 用數據說話，每點 50 字\n   - 附上數據來源\n\n3. 本月推薦物件（2-3 間）\n   - 物件亮點（不要像廣告，要像推薦）\n   - 適合什麼樣的買家\n\n4. 購屋小知識專欄\n   - 一個實用的購屋技巧（200 字）\n   - 讓讀者覺得「這個資訊值得存起來」\n\n5. CTA（行動呼籲）\n   - 引導回覆或預約看屋\n   - 不要太推銷\n\n電子報整體語氣：專業但友善，像是一位懂行的朋友在分享\n標題要吸引人打開（設計 3 個標題選項）',
    },
    {
        "num": "06",
        "title": "YouTube 長影片腳本（10分鐘）",
        "level": "🔴 進階",
        "system_note": "",
        "scenario": "你想拍一支 10 分鐘的 YouTube 影片分析板橋房市，建立長期的影片內容資產。",
        "task": "寫出完整的影片腳本",
        "tool": "ChatGPT",
        "upload_format": "📄 行情數據",
        "upload_file": "板橋區行情數據（實價登錄/新聞報導）",
        "upload_note": "上傳行情數據讓 AI 加入具體數字，增加影片可信度",
        "prompt": '我是房仲，想拍一支 10 分鐘的 YouTube 影片，主題是「2025 板橋房市完整分析」。\n\n請幫我寫完整的影片腳本：\n\n1. 前 15 秒鉤子（Hook）\n   - 用一個震撼數據或問題抓住觀眾\n   - 例如：「板橋房價今年漲了 X%，但有一種物件反而在跌…」\n\n2. 自我介紹（15 秒）\n   - 簡短有記憶點\n\n3. 正片 3 大重點（各 2-3 分鐘）\n   - 重點一：板橋各區域房價走勢比較\n   - 重點二：最值得關注的 3 個新建案\n   - 重點三：現在是不是進場的好時機？\n   - 每個重點都要有數據佐證\n\n4. 結尾 CTA（30 秒）\n   - 訂閱 + 留言互動引導\n   - 預告下一支影片\n\n額外需求：\n- 標註分鏡建議（什麼時候用什麼畫面）\n- 標註字卡內容（螢幕上要顯示的文字/數據）\n- YouTube 標題（3 個選項，要有點擊慾望）\n- 影片描述欄文字（含時間軸章節）\n- SEO 標籤（10 個）',
    },
    {
        "num": "07",
        "title": "Podcast 節目企劃+逐字稿",
        "level": "🔴 進階",
        "system_note": "",
        "scenario": "你想開一個房仲 Podcast，用對話式內容分享房市知識，但不知道怎麼規劃內容。",
        "task": "用 NotebookLM 生成對話式 Podcast + ChatGPT 寫 show notes",
        "tool": "ChatGPT + NotebookLM",
        "upload_format": "📄 房市新聞/報導",
        "upload_file": "3 篇房市新聞或分析報導",
        "upload_note": "上傳 3 篇房市新聞到 NotebookLM，讓 AI 生成對話式 Podcast",
        "prompt": '我是房仲，想開一個 Podcast 節目叫「房市內幕」。\n\n本集主題：「首購族最常犯的 5 個錯誤」\n\n請幫我設計：\n\n1. 30 分鐘節目大綱\n   - 開場白（2 分鐘）：用一個故事開場\n   - 話題一（8 分鐘）：看房前沒做功課的代價\n   - 話題二（8 分鐘）：貸款陷阱（寬限期的真相）\n   - 話題三（8 分鐘）：議價的 3 個致命錯誤\n   - 來賓 Q&A 環節（3 分鐘）：模擬聽眾提問\n   - 結尾（1 分鐘）：下集預告 + CTA\n\n2. 完整逐字稿（主持人 + 來賓對話形式）\n   - 主持人：阿峰（我）\n   - 來賓：資深代書（虛擬角色）\n   - 對話要自然，像真的在聊天\n\n3. Show Notes\n   - 本集重點摘要（5 點）\n   - 提到的資源連結\n   - 下集預告\n\n4. 社群推廣文案\n   - IG 貼文（宣傳這集 Podcast）\n   - LINE 訊息（分享給客戶群組）',
    },
    {
        "num": "08",
        "title": "SEO 部落格文章規劃",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "你想寫部落格建立自然搜尋流量，讓有購屋需求的人搜尋時找到你。",
        "task": "SEO 關鍵字研究 + 文章大綱規劃",
        "tool": "ChatGPT + Perplexity",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "先用 Perplexity 研究關鍵字搜尋量，再用 ChatGPT 規劃文章",
        "prompt": '我是板橋地區的房仲，想寫 SEO 部落格建立自然流量。\n\n請幫我：\n\n1. 關鍵字研究\n   - 「板橋買房」相關的長尾關鍵字 TOP 10\n   - 每個關鍵字標註：搜尋意圖（資訊型/交易型/導航型）\n   - 建議優先攻略的 3 個關鍵字（競爭度低但搜尋量不錯的）\n\n2. 文章大綱規劃（3 篇，每篇 2000 字）\n\n   文章一：「板橋買房完全攻略：區域、價格、交通一次看」\n   - H1 標題 + meta 描述\n   - H2/H3 子標題架構\n   - 每段的重點內容提示\n   - 內鏈建議（鏈到其他文章）\n\n   文章二：「板橋首購族必看：XXX」（你建議標題）\n   - 同上架構\n\n   文章三：「板橋 vs OO：哪裡買房CP值更高？」（你建議比較對象）\n   - 同上架構\n\n3. 內容行事曆\n   - 3 篇文章的建議發布時間\n   - 每篇文章的社群推廣計畫',
    },
    {
        "num": "09",
        "title": "社群 Q&A 問答集（批量20組）",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：社群內容批量生產系統 Batch 3",
        "scenario": "客戶常問的問題都差不多，你想一次做好 Q&A 素材，隨時可以貼到社群或存成 IG 限動。",
        "task": "批量產出 20 組 Q&A",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "不需要準備資料，直接跟 AI 對話即可",
        "prompt": '我是房仲，想一次做好社群 Q&A 素材庫。\n\n請幫我產出 20 組購屋 Q&A，要求：\n\n涵蓋以下 6 大類（每類 3-4 組）：\n1. 貸款相關（利率、成數、寬限期）\n2. 稅務相關（契稅、房地合一、土增稅）\n3. 裝修相關（預算、時間、注意事項）\n4. 議價相關（怎麼出價、斡旋技巧）\n5. 法律相關（合約、產權、糾紛）\n6. 風水相關（格局、方位、禁忌）\n\n每組 Q&A 包含：\n- 問題（用一般人會問的口語）\n- 回答（200 字，專業但易懂）\n- 一句金句（可以做成 IG 圖卡的那種）\n\n用表格呈現：編號 / 分類 / 問題 / 回答 / 金句',
    },
    {
        "num": "10",
        "title": "個人品牌定位工作坊",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "你想找到自己的品牌定位和差異化，在眾多房仲中脫穎而出。",
        "task": "用 AI 做個人品牌 SWOT 分析和定位",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "不需要準備資料，跟 AI 對話引導你做品牌定位",
        "prompt": '我是房仲，想找到自己的個人品牌定位。\n\n請用以下框架引導我做品牌定位：\n\n1. SWOT 分析\n   - 請問我以下問題，幫我整理出 SWOT：\n   - 優勢：你做得比其他房仲好的 3 件事？\n   - 劣勢：你覺得自己哪些方面還不足？\n   - 機會：市場上有什麼趨勢對你有利？\n   - 威脅：最大的競爭壓力來自哪裡？\n\n2. 目標客群定義\n   - 你最想服務什麼樣的客戶？\n   - 這群人的痛點是什麼？\n   - 他們在哪裡出沒？（線上/線下）\n\n3. 獨特賣點（USP）\n   - 根據 SWOT 結果，幫我提煉 3 個 USP\n   - 每個 USP 用一句話表達\n\n4. 品牌故事\n   - 為什麼你做房仲？\n   - 你幫客戶解決過最難忘的一個問題？\n   - 用 200 字寫出你的品牌故事\n\n5. 30 秒電梯簡報\n   - 設計一段 30 秒自我介紹\n   - 讓對方聽完會說：「我有朋友要買房，介紹給你」\n\n6. 未來 90 天行動計畫\n   - 品牌定位確定後，接下來 3 個月該做什麼？\n   - 每月 1 個重點目標 + 3 個具體行動',
    },
]

# ========== 生成案例投影片 ==========
for case in cases:
    # 情境頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 頂部標籤（含 level）
    add_text_box(slide, 0.5, 0.3, 5, 0.5, f"案例 {case['num']}  {case['level']}", 16, CYAN, True)
    add_text_box(slide, 0.5, 0.7, 12, 0.8, case["title"], 36, WHITE, True)

    # 左半邊：情境 + 任務
    tf = add_text_box(slide, 0.5, 1.8, 6, 0.5, "📋 情境", 20, ORANGE, True)
    add_paragraph(tf, case["scenario"], 17, WHITE, False, 8)
    add_paragraph(tf, "", 10, WHITE, False, 12)
    add_paragraph(tf, "🎯 任務目標", 20, ORANGE, True, 8)
    add_paragraph(tf, case["task"], 17, WHITE, False, 8)

    # 右半邊：工具 + 上傳資料
    tf2 = add_text_box(slide, 7, 1.8, 6, 0.5, "🔧 使用工具", 20, ORANGE, True)
    add_paragraph(tf2, case["tool"], 17, CYAN, True, 8)
    add_paragraph(tf2, "", 10, WHITE, False, 12)
    add_paragraph(tf2, "📎 要準備的資料", 20, ORANGE, True, 8)
    add_paragraph(tf2, f"格式：{case['upload_format']}", 15, GRAY, False, 8)
    add_paragraph(tf2, f"檔案：{case['upload_file']}", 15, WHITE, False, 4)
    add_paragraph(tf2, f"💡 {case['upload_note']}", 13, GRAY, False, 8)
    if case.get('system_note'):
        add_paragraph(tf2, case['system_note'], 13, CYAN, False, 10)

    # 操作頁（Prompt）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY)

    add_text_box(slide2, 0.5, 0.3, 6, 0.5, f"案例 {case['num']} ─ 操作指令", 16, CYAN, True)
    add_text_box(slide2, 0.5, 0.7, 12, 0.8, case["title"], 28, WHITE, True)

    # Prompt 區塊
    add_rounded_rect(slide2, 0.5, 1.6, 12.3, 5.2, CARD_BG)
    tf3 = add_text_box(slide2, 0.8, 1.7, 11.8, 0.4, "📝 複製以下指令，貼到 AI 工具中：", 16, ORANGE, True)

    prompt_lines = case["prompt"].split("\n")
    for line in prompt_lines:
        add_paragraph(tf3, line, 14, WHITE, False, 2)

    add_text_box(slide2, 0.5, 6.9, 12, 0.4, "💡 Prompt 可依實際需求修改關鍵字，不需要逐字照抄", 13, GRAY)

# ========== SOP 流程圖頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "🚀 公域獲客飛輪 SOP", 32, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.2, 12, 0.5, "每月 2 小時批量生產，永不斷更", 18, GRAY, False, PP_ALIGN.CENTER)

sop_lines = [
    "每月月初（2 小時批量生產）",
    "    ↓",
    "案例 01 → LINE 早安快訊 × 20 則（週一到週五，4 週）",
    "案例 04 → Threads 冷知識 × 10 則",
    "案例 09 → 社群 Q&A × 20 組",
    "    ↓",
    "合計 50+ 件素材，撐 4-8 週不重複",
    "    ↓",
    "每日發布：LINE 早安 + Threads（排程）",
    "每週 2-3 次：FB/IG 貼文（案例 02/03）",
    "每月一次：電子報（案例 05）",
    "    ↓",
    "進階內容：案例 06 YouTube + 案例 07 Podcast",
    "SEO 佈局：案例 08 部落格（長期流量）",
    "    ↓",
    "品牌定位：案例 10 → 所有內容統一調性",
]
tf_sop = add_text_box(slide, 1, 2.0, 11, 4.5, "", 15, WHITE)
for line in sop_lines:
    c = CYAN if "↓" in line else (ORANGE if "每月月初" in line or "合計" in line else WHITE)
    add_paragraph(tf_sop, line, 15, c, "每月月初" in line or "合計" in line, 3)

add_text_box(slide, 0.5, 6.5, 12, 0.5, "💡 社群經營最大的敵人是斷更，AI 讓你一次做好一個月的內容", 16, ORANGE, True, PP_ALIGN.CENTER)

# ========== 結尾頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 2, 11, 1, "模組 5 完成！", 44, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "你已經學會用 AI 打造公域獲客飛輪", 28, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1.5, "記住：社群經營最大的敵人是斷更\nAI 讓你一次做好一個月的內容，永不斷更\n\n💡 進階挑戰：一次跑完案例 01+04+09 = 50 件內容素材，撐 8 週", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6, 11, 0.5, "講師：阿峰老師 ｜ www.autolab.cloud", 16, CYAN, False, PP_ALIGN.CENTER)

output = "/Users/huangjingfeng/Desktop/專案/02-培訓(Training)/02-公開班(Public-Class)/03-ai-realtor-course(房仲AI課程)/cases/模組5-AI公域獲客飛輪-10案例.pptx"
prs.save(output)
print(f"完成：{output}")
print(f"共 {len(prs.slides)} 張投影片")
