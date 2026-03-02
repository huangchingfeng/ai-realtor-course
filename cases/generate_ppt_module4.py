"""模組 4：AI 追蹤成交系統 - 10 個實戰案例 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 品牌色
NAVY = RGBColor(0x0A, 0x16, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CARD_BG = RGBColor(0x0F, 0x1D, 0x32)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 0.8, 11, 1, "模組 4", 28, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "AI 追蹤成交系統", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 1, "10 個實戰案例 ─ AI 幫你追到成交的最後一哩路", 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, "AI 超級房仲實戰班 ｜ 講師：阿峰老師（黃敬峰）", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "www.autolab.cloud", 14, CYAN, False, PP_ALIGN.CENTER)

# ========== 案例資料 ==========
cases = [
    {
        "num": "01",
        "title": "帶看後會議記錄整理",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：帶看後自動跟進序列的起點（Day 0）",
        "scenario": "你剛帶客戶王大明看完板橋一間 2 房物件，過程中他提了很多想法和疑慮。你用手機錄了一段語音備忘，轉成逐字稿後想讓 AI 幫你整理重點。",
        "task": "上傳帶看逐字稿，請 AI 自動提取客戶喜好、疑慮、預算、下一步 action item",
        "tool": "ChatGPT",
        "upload_format": "📄 文字逐字稿",
        "upload_file": "帶看記錄-王大明看板橋2房.txt",
        "upload_note": "可用手機錄音後轉文字，或直接手打重點",
        "prompt": '我是房仲，剛帶客戶看完一間物件，以下是帶看過程的逐字稿。\n\n請幫我從逐字稿中提取：\n1. 客戶的喜好與在意的條件（格局、採光、交通、生活機能等）\n2. 客戶提出的疑慮或擔心的點\n3. 客戶的預算範圍和付款能力\n4. 下一步 Action Item（我該做什麼跟進動作）\n5. 客戶購買意願評估（高/中/低），並說明判斷理由\n\n請用結構化表格呈現，方便我存入客戶管理系統。',
    },
    {
        "num": "02",
        "title": "帶看後 3 天跟進訊息",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：自動跟進序列 Day 3",
        "scenario": "你 3 天前帶李美玲看了信義區一間 3 房，她當時滿喜歡的但說要跟先生討論。現在 3 天過去了，你要主動傳訊息跟進。",
        "task": "上傳帶看記錄，請 AI 生成溫暖但不過度推銷的 LINE 跟進訊息",
        "tool": "ChatGPT",
        "upload_format": "📄 文字逐字稿",
        "upload_file": "帶看記錄-李美玲看信義3房.txt",
        "upload_note": "上傳之前整理好的帶看記錄，讓 AI 根據內容客製跟進訊息",
        "prompt": '我是房仲，3 天前帶客戶李美玲看了信義區一間 3 房物件。\n以下是帶看記錄重點。\n\n請幫我寫一則 LINE 跟進訊息，要求：\n1. 語氣溫暖自然，像朋友關心而非業務推銷\n2. 提到她當天看屋時喜歡的 1-2 個亮點（從帶看記錄中提取）\n3. 自然地問她跟先生討論的結果\n4. 附帶一個有價值的資訊（例如：該社區最近的成交行情或周邊新開的店）\n5. 結尾留一個開放式問題，讓她容易回覆\n\n請產出 3 個版本：\n- 版本 A：輕鬆親切版\n- 版本 B：專業資訊版\n- 版本 C：限時推動版（暗示有其他買家在看）\n\n每則訊息控制在 100 字以內。',
    },
    {
        "num": "03",
        "title": "一週後溫度追蹤訊息",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：自動跟進序列 Day 7 / Day 14",
        "scenario": "你上週傳了跟進訊息給王大明，他只回了「好的，我再想想」就沒下文了。現在又過了一週，你不想放棄但也不想讓客戶覺得煩。",
        "task": "上傳之前的 LINE 對話紀錄，請 AI 根據對話脈絡寫出自然不尷尬的追蹤訊息",
        "tool": "ChatGPT",
        "upload_format": "📄 對話紀錄截圖或文字",
        "upload_file": "LINE對話截圖-客戶王大明.html（列印/截圖）",
        "upload_note": "將 LINE 對話截圖或複製文字上傳，讓 AI 理解前後脈絡",
        "prompt": '我是房仲，以下是我跟客戶王大明最近一週的 LINE 對話紀錄。\n\n他之前看了板橋一間 2 房，表示要考慮，上次我傳訊息他只回「好的，我再想想」。\n現在一週過去了，我想再追蹤但不想太煩人。\n\n請幫我寫追蹤訊息，要求：\n1. 不要直接問「考慮得怎麼樣」（太制式）\n2. 用一個自然的切入點開啟對話（例如：分享一個跟他需求相關的市場新訊息）\n3. 讓他感覺你是在提供價值，不是在催他\n4. 如果他已讀不回，再給我一個 3 天後的第二波追蹤訊息\n\n產出：\n- 第一波訊息（今天發）\n- 第二波訊息（3 天後發，如果第一波沒回）\n- 第三波訊息（一週後發，最後一次溫和追蹤）',
    },
    {
        "num": "04",
        "title": "客戶異議處理 —「太貴了」",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "李美玲看完信義區 3 房後說：「這個價格太貴了，隔壁社區便宜好幾百萬。」你需要用數據來回應她的異議。",
        "task": "上傳物件資料與實價登錄，請 AI 用數據回應「太貴」的異議，產出 3 種不同語氣的回覆",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案 + 📄 物件資料",
        "upload_file": "物件資料+實價登錄CSV",
        "upload_note": "上傳物件開價資料和同區實價登錄，讓 AI 用數據說話",
        "prompt": '我是房仲，客戶看完信義區一間 3 房後說「太貴了」。\n\n物件開價：4,200 萬，坪數 45 坪（含車位）\n以下是附近社區的實價登錄資料。\n\n請幫我：\n1. 分析這間物件的開價是否合理（跟同社區、同區域比較）\n2. 找出這間物件比隔壁社區「貴得有道理」的點（屋齡、格局、管理、建材等）\n3. 產出 3 種不同語氣的回覆：\n   - 版本 A：溫和理解型（先同理再解釋）\n   - 版本 B：數據說服型（直接用數字比較）\n   - 版本 C：價值重塑型（轉移焦點到生活品質和長期增值）\n\n每個版本都要引用具體數據，不能只講空話。\n控制在 LINE 訊息的長度（150 字以內）。',
    },
    {
        "num": "05",
        "title": "客戶異議處理 —「再考慮看看」",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "你帶了趙小萱看了 3 間房，每次她都說「還不錯，但我再考慮看看」。你感覺她有興趣但總是下不了決定。",
        "task": "請 AI 分析客戶猶豫的可能原因，並給出 5 種不同的跟進策略",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個案例不需要準備資料，直接跟 AI 對話即可",
        "prompt": '我是房仲，有一位客戶趙小萱，首購族，28歲，預算約 1,000 萬。\n\n我帶她看了 3 間房子：\n- A 物件：板橋套房 15 坪，850 萬（她說格局不錯但太小）\n- B 物件：三重 2 房 25 坪，980 萬（她說地點不夠好）\n- C 物件：中和 2 房 22 坪，1,050 萬（她說超出預算一點）\n\n每次看完她都說「還不錯，我再考慮看看」，已經拖了 2 週沒決定。\n\n請幫我分析：\n1. 她可能猶豫的 5 個原因（心理層面分析）\n2. 針對每個原因，給出對應的跟進策略\n3. 一則 LINE 訊息，用「不推銷」的方式幫她釐清真正的需求\n4. 如果要約她再看一次，怎麼說比較不會被拒絕？\n5. 判斷：她到底是真的要買還是只是在逛？有什麼方法測試？',
    },
    {
        "num": "06",
        "title": "多客戶同時管理",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：CRM 溫度排序，每週一自動產出優先聯繫榜",
        "scenario": "你手上同時有 10 位客戶在跟，每個人的需求、進度、溫度都不一樣，你快搞不清楚該先聯絡誰。",
        "task": "上傳客戶名單，請 AI 根據溫度排序，產出本週優先聯繫順序和各客戶的跟進訊息",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案",
        "upload_file": "客戶名單與需求.csv",
        "upload_note": "CSV 欄位建議：姓名、需求、預算、上次聯繫日、溫度、備註",
        "prompt": '我是房仲，以下是我目前手上 10 位客戶的資料（CSV 檔）。\n\n請幫我：\n1. 根據以下因素進行「客戶溫度排序」：\n   - 上次聯繫距今天數（越久越需要聯繫）\n   - 購買意願（高/中/低）\n   - 預算與市場匹配度（預算合理的優先）\n   - 客戶類型（投資客通常決策較快）\n\n2. 產出「本週優先聯繫排行榜」（前 5 名）\n\n3. 針對排行榜上的每位客戶，各寫一則客製化的 LINE 跟進訊息\n   （要根據該客戶的需求和上次互動內容來寫）\n\n4. 對於溫度較低的客戶，建議是「繼續跟進」還是「暫時放下」\n\n請用表格呈現排序結果，訊息分開列出。',
    },
    {
        "num": "07",
        "title": "斡旋金談判話術",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "買方王大明出價 1,500 萬，賣方底價 1,600 萬，差距 100 萬。你夾在中間，需要兩邊都說服。",
        "task": "請 AI 產出買方端和賣方端的談判策略與話術",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個案例練習談判策略，直接跟 AI 對話即可",
        "prompt": '我是房仲，目前有一個斡旋案：\n\n物件：板橋新板特區 2 房，32 坪，屋齡 8 年\n- 賣方開價：1,680 萬\n- 賣方底價：1,600 萬（我私下知道的）\n- 買方出價：1,500 萬（已下斡旋金 10 萬）\n- 價差：100 萬\n\n買方狀況：首購族，預算有限，但真的很喜歡這間\n賣方狀況：急售（要換屋），但心理上覺得賠太多不甘心\n\n請幫我產出完整的談判策略：\n\n1. 對買方的話術（讓他願意加價到 1,550-1,580）\n   - 3 個說服理由\n   - 具體話術範例\n\n2. 對賣方的話術（讓他願意降到 1,580-1,600）\n   - 3 個說服理由\n   - 具體話術範例\n\n3. 斡旋談判 SOP：先找誰談？每一步怎麼做？\n4. 如果談不攏的 Plan B 是什麼？',
    },
    {
        "num": "08",
        "title": "已讀不回處理 SOP",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：自動偵測沉默天數 → 觸發對應話術模板",
        "scenario": "你手上有 3 位客戶同時已讀不回：王大明看完房子 7 天沒回、李美玲 3 週沒消息、趙小萱上個月說要考慮後就消失了。每個人的沉默程度不同，你不能用同一套話術。你需要一個系統性的處理 SOP，根據沉默天數決定該說什麼。",
        "task": "請 AI 根據「沉默天數」產出分級處理策略，從 7 天到 90 天各有不同的話術和行動方案",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個 SOP 做好後可以存成模板，每次遇到已讀不回直接套用",
        "prompt": '我是房仲，目前有多位客戶「已讀不回」，我需要一套系統化的處理 SOP。\n\n請根據「沉默天數」幫我設計分級處理策略：\n\n第一級：沉默 3-7 天\n- 可能原因分析（3 個最常見原因）\n- LINE 訊息模板（2 個版本）\n- 切入角度：提供新價值，不問結果\n\n第二級：沉默 7-14 天\n- 可能原因分析\n- LINE 訊息模板（2 個版本）\n- 切入角度：用一個跟他需求相關的市場新消息重啟對話\n- 備選方案：打電話還是傳訊息？\n\n第三級：沉默 14-30 天\n- 可能原因分析\n- LINE 訊息模板（2 個版本）\n- 切入角度：降低壓力，給他一個「不買也沒關係」的台階\n- 是否應該直接問：「是不是已經在其他地方看到喜歡的了？」\n\n第四級：沉默 30-90 天\n- 判斷：這個客戶是放棄還是暫緩？\n- 「最後一搏」訊息模板\n- 何時應該從「主動跟進」改為「被動維繫」（放到節日關懷名單）\n\n第五級：沉默超過 90 天\n- 降級到「養客名單」的處理方式\n- 每季一次的「價值投放」訊息模板\n\n另外，請幫我做一張「已讀不回處理決策樹」：\n已讀不回 → 幾天了？ → 對應動作 → 成功/繼續沉默 → 下一步\n\n每個模板都要自然、不尷尬、不讓客戶覺得被施壓。',
    },
    {
        "num": "09",
        "title": "成交後感謝訊息+轉介紹請求",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：成交後 4 節點自動序列（Day0→Day7→Day30→Day90）",
        "scenario": "恭喜！王大明的案子成交了！但成交不是結束，而是經營客戶關係的開始。你要在不同時間點傳送感謝訊息，並自然地請求轉介紹。",
        "task": "請 AI 產出 3 個不同時間點的訊息：成交當天、交屋後一週、一個月後",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個案例練習長期客戶關係經營，直接跟 AI 對話",
        "prompt": '我是房仲，客戶王大明剛買了板橋新板特區一間 2 房，成交價 1,580 萬。\n他是首購族，32 歲，跟女友一起住，之後打算結婚。\n整個過程他對我的服務很滿意。\n\n請幫我寫 3 個時間點的 LINE 訊息：\n\n1. 成交當天（感謝 + 恭喜）\n   - 真誠感謝，不要太制式\n   - 提醒接下來的交屋流程重點\n\n2. 交屋後一週（關心入住 + 小禮物）\n   - 問入住狀況，有沒有需要幫忙的\n   - 推薦附近好吃的店、生活便利資訊\n   - 如果要送小禮物，建議送什麼？\n\n3. 一個月後（轉介紹請求）\n   - 自然地帶到「如果身邊有朋友要買房…」\n   - 不要太刻意，要讓他覺得是順便提\n   - 附帶一個「推薦好友送禮券」的機制設計\n\n每則訊息控制在 120 字以內，語氣像朋友不像業務。\n\n額外請求：\n6. 第四個時間點：成交 90 天後（二次轉介紹 + 市場近況更新）\n   - 距離上次轉介紹已經 2 個月，自然地二次提醒\n   - 順便分享他房子附近的最新成交行情（讓他知道房價漲了，心情好更願意推薦）\n   - 控制在 100 字以內\n\n7. 設計「轉介紹禮遇卡」的文案\n   - 介紹朋友成功成交，送什麼？怎麼設計最吸引人？\n   - 一張可以轉發的 LINE 圖卡文案',
    },
    {
        "num": "10",
        "title": "客戶生日/節日關懷訊息",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：全年客戶關懷行事曆，節日前 2 天自動提醒",
        "scenario": "你手上有 10 位客戶，有的是首購族、有的是投資客、有的是換屋族。農曆新年快到了，你想傳關懷訊息維繫關係。",
        "task": "上傳客戶名單，請 AI 根據客戶類型產出不同風格的關懷訊息模板",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案",
        "upload_file": "客戶名單與需求.csv",
        "upload_note": "CSV 中要有客戶類型欄位（首購/投資/換屋），讓 AI 客製化訊息",
        "prompt": '我是房仲，以下是我的客戶名單（CSV 檔），農曆新年快到了。\n\n請根據不同客戶類型，幫我產出關懷訊息模板：\n\n1. 首購族客戶（年輕、第一次買房）\n   - 語氣：活潑、像朋友\n   - 內容：新年快樂 + 新家開運小 tips\n\n2. 投資客客戶（注重數字和效率）\n   - 語氣：專業簡潔\n   - 內容：新年快樂 + 新年度房市展望一句話\n\n3. 換屋族客戶（家庭取向、注重生活品質）\n   - 語氣：溫暖關懷\n   - 內容：新年快樂 + 新家/新生活的祝福\n\n4. 已成交客戶（維繫關係、請求轉介紹）\n   - 語氣：感恩回饋\n   - 內容：感謝信任 + 新年小禮物/優惠\n\n5. 尚未成交但持續關注的客戶\n   - 語氣：不推銷、提供價值\n   - 內容：新年快樂 + 分享一個有用的房市資訊\n\n每則訊息 80 字以內，要能直接複製貼上到 LINE。\n另外幫我做一張「全年客戶關懷行事曆」（包含：農曆新年、端午、中秋、生日、成交周年等時間點）。',
    },
]

# ========== 生成案例投影片 ==========
for case in cases:
    # 情境頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 頂部標籤（含 level）
    add_text_box(slide, 0.5, 0.3, 5, 0.5, f"案例 {case['num']}  {case['level']}", 16, CYAN, True)
    add_text_box(slide, 0.5, 0.7, 12, 0.8, case["title"], 36, WHITE, True)

    # 左半邊：情境 + 任務
    tf = add_text_box(slide, 0.5, 1.8, 6, 0.5, "📋 情境", 20, ORANGE, True)
    add_paragraph(tf, case["scenario"], 17, WHITE, False, 8)
    add_paragraph(tf, "", 10, WHITE, False, 12)
    add_paragraph(tf, "🎯 任務目標", 20, ORANGE, True, 8)
    add_paragraph(tf, case["task"], 17, WHITE, False, 8)

    # 右半邊：工具 + 上傳資料
    tf2 = add_text_box(slide, 7, 1.8, 6, 0.5, "🔧 使用工具", 20, ORANGE, True)
    add_paragraph(tf2, case["tool"], 17, CYAN, True, 8)
    add_paragraph(tf2, "", 10, WHITE, False, 12)
    add_paragraph(tf2, "📎 要準備的資料", 20, ORANGE, True, 8)
    add_paragraph(tf2, f"格式：{case['upload_format']}", 15, GRAY, False, 8)
    add_paragraph(tf2, f"檔案：{case['upload_file']}", 15, WHITE, False, 4)
    add_paragraph(tf2, f"💡 {case['upload_note']}", 13, GRAY, False, 8)
    if case.get('system_note'):
        add_paragraph(tf2, case['system_note'], 13, CYAN, False, 10)

    # 操作頁（Prompt）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY)

    add_text_box(slide2, 0.5, 0.3, 6, 0.5, f"案例 {case['num']} ─ 操作指令", 16, CYAN, True)
    add_text_box(slide2, 0.5, 0.7, 12, 0.8, case["title"], 28, WHITE, True)

    # Prompt 區塊
    add_rounded_rect(slide2, 0.5, 1.6, 12.3, 5.2, CARD_BG)
    tf3 = add_text_box(slide2, 0.8, 1.7, 11.8, 0.4, "📝 複製以下指令，貼到 AI 工具中：", 16, ORANGE, True)

    prompt_lines = case["prompt"].split("\n")
    for line in prompt_lines:
        add_paragraph(tf3, line, 14, WHITE, False, 2)

    add_text_box(slide2, 0.5, 6.9, 12, 0.4, "💡 Prompt 可依實際物件修改關鍵字，不需要逐字照抄", 13, GRAY)

# ========== SOP 流程圖頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "客戶追蹤成交 SOP", 32, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.2, 12, 0.5, "從帶看到成交到轉介紹，每個節點都有 AI 幫你", 18, GRAY, False, PP_ALIGN.CENTER)

sop_lines = [
    "帶看結束 → 案例 01 逐字稿自動整理 → 存入客戶資料卡",
    "    ↓",
    "Day 3 → 案例 02 自動產出跟進訊息（3 版本選 1）",
    "    ↓",
    "Day 7 → 案例 03 溫度追蹤訊息",
    "    ↓  （如果客戶已讀不回）→ 案例 08 分級處理 SOP",
    "    ↓",
    "客戶異議 → 案例 04「太貴了」/ 案例 05「再考慮」",
    "    ↓",
    "進入斡旋 → 案例 07 談判策略（買方端 + 賣方端）",
    "    ↓",
    "成交！ → 案例 09 四節點自動序列（Day0→7→30→90）",
    "    ↓",
    "長期維繫 → 案例 10 全年關懷行事曆 → 轉介紹飛輪",
    "",
    "每週一早上 → 案例 06 CRM 溫度排序 → 今日優先聯繫榜",
]
tf_sop = add_text_box(slide, 1, 2.0, 11, 4.5, "", 15, WHITE)
for line in sop_lines:
    c = CYAN if "↓" in line else (ORANGE if "每週一" in line else WHITE)
    add_paragraph(tf_sop, line, 15, c, "每週一" in line, 3)

add_text_box(slide, 0.5, 6.5, 12, 0.5, "💡 成交不是靠運氣，是靠系統化跟進。AI 幫你記住每位客戶", 16, ORANGE, True, PP_ALIGN.CENTER)

# ========== 結尾頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 2, 11, 1, "模組 4 完成！", 44, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "你已經學會用 AI 建立追蹤成交系統", 28, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1.5, "記住：成交不是靠運氣，是靠系統化跟進\nAI 幫你記住每位客戶，你負責用心經營關係\n\n💡 進階挑戰：把案例 01→02→03→08 串成「帶看後自動跟進序列」", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6, 11, 0.5, "講師：阿峰老師 ｜ www.autolab.cloud", 16, CYAN, False, PP_ALIGN.CENTER)

output = "/Users/huangjingfeng/Desktop/專案/02-培訓(Training)/02-公開班(Public-Class)/03-ai-realtor-course(房仲AI課程)/cases/模組4-AI追蹤成交系統-10案例.pptx"
prs.save(output)
print(f"完成：{output}")
print(f"共 {len(prs.slides)} 張投影片")
