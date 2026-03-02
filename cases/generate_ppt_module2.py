"""模組 2：AI 簡報製作 - 10 個實戰案例 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 品牌色
NAVY = RGBColor(0x0A, 0x16, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CARD_BG = RGBColor(0x0F, 0x1D, 0x32)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 0.8, 11, 1, "模組 2", 28, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "📑 AI 簡報製作", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 1, "10 個實戰案例 ─ 用 AI 做出專業簡報", 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, "AI 超級房仲實戰班 ｜ 講師：阿峰老師（黃敬峰）", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "www.autolab.cloud", 14, CYAN, False, PP_ALIGN.CENTER)

# ========== 案例資料 ==========
cases = [
    {
        "num": "01",
        "title": "社區介紹簡報",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "你要介紹信義區翠亨村社區給一位新客戶，想做一份專業簡報，讓客戶快速了解社區的優勢和生活機能。",
        "task": "用 AI 整理社區資訊，再用 Gamma 生成一份專業的社區介紹 PPT",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 社區照片 + 基本資料",
        "upload_file": "社區照片+基本資料文件",
        "upload_note": "準備社區外觀、公設、周邊環境的照片，以及管委會提供的基本資料",
        "prompt": '我是房仲，要做一份「信義區翠亨村社區」的介紹簡報給新客戶看。\n\n請幫我整理以下資訊，並以 12 頁簡報大綱呈現：\n\n第 1 頁：封面（社區名稱 + 一句話定位）\n第 2 頁：社區基本資料（地址/屋齡/戶數/坪數範圍/車位）\n第 3 頁：管委會與社區管理（管理費/保全/清潔/管委會運作）\n第 4 頁：公共設施介紹（健身房/游泳池/中庭花園/閱覽室等）\n第 5 頁：格局與坪數分析（主力格局/各房型比例）\n第 6 頁：近一年實價登錄行情（均價/最高/最低/趨勢）\n第 7 頁：周邊生活機能 - 交通（捷運/公車/開車動線）\n第 8 頁：周邊生活機能 - 學區與公園\n第 9 頁：周邊生活機能 - 商圈與餐飲\n第 10 頁：社區優勢總結（3 大賣點）\n第 11 頁：適合什麼樣的買家？（TA 分析）\n第 12 頁：聯絡資訊與下一步\n\n每頁請列出標題和 3-5 個重點內容，方便我直接貼到 Gamma 生成簡報。',
    },
    {
        "num": "02",
        "title": "客戶提案簡報（個人化方案）",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：帶看前自動產出客戶提案",
        "scenario": "你要幫客戶王大明量身打造一份購屋提案，根據他的預算、需求和通勤路線，推薦最適合的物件。",
        "task": "整合客戶需求、行情數據和推薦物件，做成個人化的購屋提案簡報",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 客戶需求表 + 推薦物件清單",
        "upload_file": "客戶需求表-王大明.txt + 推薦物件清單.csv",
        "upload_note": "準備客戶的需求（預算/坪數/格局/通勤地點），以及你篩選出的 3-5 間推薦物件",
        "prompt": '我是房仲，要幫客戶王大明做一份個人化的購屋提案簡報。\n\n客戶資料：\n- 姓名：王大明，32 歲，科技業工程師\n- 預算：1,200-1,500 萬\n- 需求：2 房，20-30 坪，要有車位\n- 通勤：南港軟體園區（希望捷運 30 分鐘內）\n- 偏好：屋齡 10 年內，有管理的社區\n\n請根據以上需求，產出個人化購屋提案簡報大綱：\n\n第 1 頁：封面（王大明專屬購屋提案）\n第 2 頁：需求摘要（用表格整理他的需求條件）\n第 3 頁：區域推薦分析（哪些區域符合通勤+預算條件）\n第 4 頁：推薦物件 A（基本資料/優點/注意事項）\n第 5 頁：推薦物件 B\n第 6 頁：推薦物件 C\n第 7 頁：三物件比較表（價格/坪數/屋齡/通勤時間/管理費）\n第 8 頁：行情分析（推薦區域的實價登錄趨勢）\n第 9 頁：購屋成本試算（頭期/月付/車位/管理費）\n第 10 頁：看屋行程建議（建議帶看順序和時間）\n第 11 頁：為什麼選我當你的經紀人？\n第 12 頁：下一步行動 + 聯絡方式\n\n每頁列出標題和 3-5 個重點，語氣專業但親切。',
    },
    {
        "num": "03",
        "title": "區域月報簡報",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：區域月報自動生成的簡報版",
        "scenario": "每個月你都要做一份信義區的行情報告給客戶，讓他們掌握市場動態。以前用 Excel 拉圖表很花時間，現在用 AI 來做。",
        "task": "上傳實價登錄 CSV 資料，請 AI 分析後產出月報簡報大綱",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 CSV 實價登錄資料",
        "upload_file": "信義區實價登錄-2024年.csv",
        "upload_note": "從實價登錄網站下載該區域的 CSV 資料，上傳給 ChatGPT 分析",
        "prompt": '我是房仲，以下是信義區最近一個月的實價登錄資料（CSV 檔）。\n\n請幫我分析資料，並產出 12 頁的區域月報簡報大綱：\n\n第 1 頁：封面（信義區不動產月報 - 2024 年 X 月）\n第 2 頁：本月重點摘要（3 個關鍵數字）\n第 3 頁：成交量分析（本月 vs 上月 vs 去年同期）\n第 4 頁：均價走勢（近 6 個月折線圖的數據）\n第 5 頁：各坪數段成交分布（小宅/2房/3房/豪宅）\n第 6 頁：TOP 5 成交社區排行\n第 7 頁：單價最高 vs 最低案例分析\n第 8 頁：新建案 vs 中古屋比較\n第 9 頁：租金行情與投報率\n第 10 頁：周邊區域比較（信義 vs 大安 vs 松山）\n第 11 頁：市場趨勢解讀與投資建議\n第 12 頁：聯絡我取得完整報告\n\n每頁列出標題、關鍵數據和 2-3 個分析觀點。\n數據請從我上傳的 CSV 中提取，不要用假數據。',
    },
    {
        "num": "04",
        "title": "開發簡報（屋主委售提案）",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "一位屋主想賣房子，但同時有 3 個房仲在競爭。你需要做一份專業的委售提案，說服屋主選你當他的經紀人。",
        "task": "做一份專業的委售提案簡報，展現你的行銷計畫和銷售策略",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 物件照片 + 該區行情資料",
        "upload_file": "物件照片+該區實價登錄行情",
        "upload_note": "準備物件的照片、基本資料，以及同區域的成交行情作為定價參考",
        "prompt": '我是房仲，有一位屋主想賣信義區的 3 房物件（45 坪，屋齡 12 年）。\n目前有 3 個房仲在競爭這個委售案，我要做一份最專業的提案勝出。\n\n請幫我產出委售提案簡報大綱：\n\n第 1 頁：封面（專屬銷售提案 - 信義區 XX 路 X 號）\n第 2 頁：物件亮點分析（這間房子的 5 大優勢）\n第 3 頁：市場行情分析（同社區/同區域近半年成交行情）\n第 4 頁：建議售價策略（開價/底價/預期成交價，含理由）\n第 5 頁：行銷計畫 - 線上（591/社群/Google 廣告投放）\n第 6 頁：行銷計畫 - 線下（帶看安排/Open House/DM）\n第 7 頁：專業攝影與文案（拍攝計畫/文案風格預覽）\n第 8 頁：銷售時程表（上架→帶看→斡旋→成交，預估 60 天）\n第 9 頁：我的成交實績（同區域/同類型物件的銷售紀錄）\n第 10 頁：服務承諾（每週回報/帶看紀錄透明化/不灌水帶看）\n第 11 頁：為什麼選我？（vs 其他房仲的差異化優勢）\n第 12 頁：下一步：簽約委售 + 聯絡方式\n\n語氣要專業自信，用數據說話，讓屋主覺得「這個人很不一樣」。',
    },
    {
        "num": "05",
        "title": "投資分析簡報",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "投資客張先生要你做一份完整的投資分析，他想評估一間收租套房的投資報酬率。數字要精準，不能只講感覺。",
        "task": "做成一份投資報酬率分析簡報，含購入成本、租金收益、毛淨報酬率",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 物件資料 + 租金行情",
        "upload_file": "物件資料+該區租金行情.csv",
        "upload_note": "準備物件的售價、坪數、屋齡，以及同區域的租金行情資料",
        "prompt": '我是房仲，投資客張先生想評估一間收租套房的投資價值。\n\n物件資料：\n- 位置：中山區南京東路，捷運步行 3 分鐘\n- 坪數：15 坪（權狀），室內約 12 坪\n- 屋齡：8 年\n- 售價：880 萬\n- 目前租金行情：月租約 18,000-22,000\n- 管理費：每月 2,500\n\n請幫我產出投資分析簡報大綱：\n\n第 1 頁：封面（投資分析報告 - 中山區南京東路套房）\n第 2 頁：物件基本資料總覽\n第 3 頁：購入成本試算（房價+契稅+仲介費+裝潢/整理費）\n第 4 頁：租金收益分析（月租/年租/空置率估算）\n第 5 頁：毛投報率計算（年租金÷總購入成本）\n第 6 頁：淨投報率計算（扣除管理費/稅/維修/空置後）\n第 7 頁：5 年持有模擬（租金累計+房價增值預估）\n第 8 頁：10 年持有模擬（含房貸還款進度）\n第 9 頁：風險評估（利率上升/空置/房價下跌情境）\n第 10 頁：同區域比較（vs 其他 3 間類似物件的投報率）\n第 11 頁：增值潛力分析（捷運/商圈/都更/重大建設）\n第 12 頁：投資建議與下一步\n\n所有數字請精確計算，不要用模糊的描述。投報率請算到小數點第二位。',
    },
    {
        "num": "06",
        "title": "建案比較簡報",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "客戶在看 3 個新建案，每個建案的 DM 都說自己最好。你要做一份客觀的比較表，幫客戶做出最好的選擇。",
        "task": "上傳 3 份建案 DM 或文宣，請 AI 做出客觀的比較分析簡報",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 3 份建案 DM / 文宣",
        "upload_file": "建案A-DM.pdf + 建案B-DM.pdf + 建案C-DM.pdf",
        "upload_note": "上傳 3 份建案的 DM、文宣或官網截圖，讓 AI 提取關鍵資訊做比較",
        "prompt": '我是房仲，客戶正在考慮 3 個新建案，需要我做一份客觀的比較分析。\n\n以下是 3 個建案的 DM 資料（已上傳）。\n\n請幫我產出建案比較簡報大綱：\n\n第 1 頁：封面（三建案比較分析 - 幫你找到最適合的家）\n第 2 頁：三建案基本資料總覽表（建案名/建商/位置/坪數/單價/總價）\n第 3 頁：格局比較（各建案的主力格局/實際使用坪數/公設比）\n第 4 頁：建材與規格比較（結構/外觀/廚具/衛浴/地板）\n第 5 頁：公設與管理比較（公設項目/管理費/保全/車位）\n第 6 頁：地段與交通比較（捷運距離/生活機能/學區）\n第 7 頁：價格與付款比較（單價/總價/付款方式/優惠）\n第 8 頁：建商信譽比較（過去建案品質/交屋評價/財務狀況）\n第 9 頁：優缺點總結（每個建案 3 優點 + 2 注意事項）\n第 10 頁：不同需求推薦（自住推薦/投資推薦/首購推薦）\n第 11 頁：預計交屋時程與注意事項\n第 12 頁：我的建議 + 下一步看屋安排\n\n比較要客觀中立，列出數據，不偏袒任何一個建案。',
    },
    {
        "num": "07",
        "title": "店面投資提案",
        "level": "🔴 進階",
        "system_note": "",
        "scenario": "投資客想買東區一間店面，商用不動產的分析比住宅複雜得多，需要考慮商圈人流、租金水平、業種分析等。",
        "task": "做一份商用不動產投資分析簡報，含商圈分析和投報率計算",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 店面資料 + 商圈分析資料",
        "upload_file": "店面資料+東區商圈分析.pdf",
        "upload_note": "準備店面的基本資料、租金行情，以及商圈的人流和業種分布資料",
        "prompt": '我是房仲，投資客想買東區忠孝東路四段一間店面。\n\n店面資料：\n- 位置：忠孝東路四段，一樓店面\n- 坪數：25 坪（前面 15 坪店面 + 後面 10 坪倉庫）\n- 售價：6,800 萬\n- 目前出租中，月租 15 萬（租約剩 2 年）\n- 目前承租業種：服飾店\n\n請幫我產出店面投資分析簡報大綱：\n\n第 1 頁：封面（東區店面投資分析報告）\n第 2 頁：物件基本資料與現況\n第 3 頁：商圈分析 - 東區商圈定位與發展趨勢\n第 4 頁：人流分析（平日/假日/尖峰時段預估）\n第 5 頁：周邊業種分布（餐飲/服飾/美妝/其他比例）\n第 6 頁：租金行情比較（同路段/同坪數的租金水平）\n第 7 頁：投報率計算（毛投報率/淨投報率/含房貸試算）\n第 8 頁：租約分析與租約到期後的策略\n第 9 頁：風險評估（商圈轉移/電商衝擊/租客違約）\n第 10 頁：同區域成交案例比較（近 2 年同路段店面成交）\n第 11 頁：增值潛力（東區復甦趨勢/捷運TOD/都更）\n第 12 頁：投資建議與下一步\n\n商用不動產分析要比住宅更嚴謹，所有數字要有依據。',
    },
    {
        "num": "08",
        "title": "公司形象簡報",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "你加入新的房仲品牌，要做一份個人形象簡報。不管是開發客戶、參加社區說明會，還是自我介紹，都能用這份簡報。",
        "task": "做一份個人品牌介紹 PPT，展現你的專業度和服務理念",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 個人照片 + 經歷資料",
        "upload_file": "個人照片+經歷整理.txt",
        "upload_note": "準備個人專業照、經歷、成交實績、客戶推薦語等素材",
        "prompt": '我是房仲，剛加入信義房屋，要做一份個人形象簡報。\n\n我的基本資料：\n- 姓名：陳小明，35 歲\n- 經歷：從業 5 年，之前在永慶\n- 專精區域：信義區、大安區\n- 成交實績：累計成交 120 件\n- 專長：豪宅市場、換屋族服務\n- 座右銘：「用心，讓每個家都找到對的人」\n\n請幫我產出個人品牌簡報大綱：\n\n第 1 頁：封面（個人形象照 + 姓名 + 一句話定位）\n第 2 頁：自我介紹（為什麼選擇做房仲？我的初心）\n第 3 頁：服務理念（3 大服務承諾）\n第 4 頁：專業能力（市場分析/談判/行銷的具體能力）\n第 5 頁：專精區域介紹（我最熟的 2 個區域）\n第 6 頁：成交實績數字（120 件成交/服務 X 組客戶/平均成交天數）\n第 7 頁：代表性成交案例 1（含故事）\n第 8 頁：代表性成交案例 2（含故事）\n第 9 頁：客戶推薦語（3-5 則真實推薦）\n第 10 頁：我能幫你什麼？（針對不同客戶類型的服務）\n第 11 頁：服務流程（從初次諮詢到成交交屋的完整流程）\n第 12 頁：聯絡方式 + CTA（掃 QR code 加 LINE）\n\n語氣要真誠專業，不要太浮誇，讓客戶感覺可以信任。',
    },
    {
        "num": "09",
        "title": "成交報告簡報",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：成交後自動觸發，用於客戶分享和轉介紹",
        "scenario": "王大明的案子成交了！你要做一份成交報告，既是給客戶的恭喜和交屋指南，也是未來轉介紹的素材。",
        "task": "做一份成交恭喜簡報，包含交屋流程和轉介紹引導",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 物件照片 + 成交資訊",
        "upload_file": "成交物件照片+成交資訊.txt",
        "upload_note": "準備成交物件的照片、成交價格、交屋日期等資訊",
        "prompt": '我是房仲，客戶王大明剛成交了板橋新板特區一間 2 房物件。\n\n成交資訊：\n- 物件：板橋新板特區 XX 社區 2 房\n- 成交價：1,580 萬\n- 坪數：28 坪（含車位）\n- 交屋日：2024 年 3 月 15 日\n- 客戶背景：首購族，32 歲，跟女友一起住\n\n請幫我產出成交報告簡報大綱：\n\n第 1 頁：封面（恭喜王大明，新家落成！）\n第 2 頁：物件亮點回顧（這間房子的 5 大優勢）\n第 3 頁：成交歷程回顧（從看屋到成交的時間軸）\n第 4 頁：交屋流程說明（驗屋→過戶→交屋，每步驟注意事項）\n第 5 頁：交屋檢查清單（水電/門窗/地板/廚衛等 check list）\n第 6 頁：入住前準備（水電過戶/網路申請/管委會報到）\n第 7 頁：周邊生活指南（推薦餐廳/超市/醫院/公園）\n第 8 頁：新家布置建議（根據坪數推薦的風格和預算）\n第 9 頁：房屋保值小知識（維護保養/社區參與/增值改造）\n第 10 頁：感謝頁（感謝信任，這是我們一起完成的）\n第 11 頁：服務延續（未來有任何房屋問題都可以找我）\n第 12 頁：轉介紹頁（分享給身邊有買房需求的朋友）\n\n語氣溫暖恭喜，讓客戶覺得被重視，自然願意轉介紹。',
    },
    {
        "num": "10",
        "title": "團隊季度報告",
        "level": "🔴 進階",
        "system_note": "",
        "scenario": "你是店長，要做一份 Q1 季報給總公司。數據很多、面向很廣，手動做要花一整天。用 AI 幫你整理數據、產出簡報。",
        "task": "上傳團隊業績表 Excel，做成一份團隊業績季報簡報",
        "tool": "ChatGPT + Gamma",
        "upload_format": "📄 團隊業績表 Excel",
        "upload_file": "Q1團隊業績表.xlsx",
        "upload_note": "上傳包含每位業務員的業績、成交件數、客戶數等數據的 Excel",
        "prompt": '我是房仲店長，以下是我的團隊 Q1 業績表（Excel 檔）。\n\n團隊資料：\n- 店名：信義房屋信義旗艦店\n- 人數：12 位業務\n- Q1 期間：2024 年 1-3 月\n\n請幫我產出 Q1 季度報告簡報大綱：\n\n第 1 頁：封面（信義旗艦店 Q1 季度報告）\n第 2 頁：Q1 業績總覽（總成交金額/件數/達成率，用大數字呈現）\n第 3 頁：業績達成率分析（目標 vs 實際，各月趨勢）\n第 4 頁：人均產能分析（平均每人成交件數/金額/佣金）\n第 5 頁：TOP 3 業務表揚（成交金額/件數/客戶滿意度排行）\n第 6 頁：各業務員業績明細表\n第 7 頁：客戶來源分析（網路/推薦/開發/舊客各佔比）\n第 8 頁：成交物件分析（區域/坪數/價格帶分布）\n第 9 頁：客戶滿意度統計（NPS/回饋摘要）\n第 10 頁：Q1 問題與檢討（哪些地方可以改善）\n第 11 頁：Q2 目標與策略（業績目標/重點推動事項）\n第 12 頁：團隊合照 + 一句話激勵\n\n數據請從 Excel 中提取，圖表建議用哪種呈現方式也請標註。',
    },
]

# ========== 生成案例投影片 ==========
for case in cases:
    # 情境頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 頂部標籤（含 level）
    add_text_box(slide, 0.5, 0.3, 5, 0.5, f"案例 {case['num']}  {case['level']}", 16, CYAN, True)
    add_text_box(slide, 0.5, 0.7, 12, 0.8, case["title"], 36, WHITE, True)

    # 左半邊：情境 + 任務
    tf = add_text_box(slide, 0.5, 1.8, 6, 0.5, "📋 情境", 20, ORANGE, True)
    add_paragraph(tf, case["scenario"], 17, WHITE, False, 8)
    add_paragraph(tf, "", 10, WHITE, False, 12)
    add_paragraph(tf, "🎯 任務目標", 20, ORANGE, True, 8)
    add_paragraph(tf, case["task"], 17, WHITE, False, 8)

    # 右半邊：工具 + 上傳資料
    tf2 = add_text_box(slide, 7, 1.8, 6, 0.5, "🔧 使用工具", 20, ORANGE, True)
    add_paragraph(tf2, case["tool"], 17, CYAN, True, 8)
    add_paragraph(tf2, "", 10, WHITE, False, 12)
    add_paragraph(tf2, "📎 要準備的資料", 20, ORANGE, True, 8)
    add_paragraph(tf2, f"格式：{case['upload_format']}", 15, GRAY, False, 8)
    add_paragraph(tf2, f"檔案：{case['upload_file']}", 15, WHITE, False, 4)
    add_paragraph(tf2, f"💡 {case['upload_note']}", 13, GRAY, False, 8)
    if case.get('system_note'):
        add_paragraph(tf2, case['system_note'], 13, CYAN, False, 10)

    # 操作頁（Prompt）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY)

    add_text_box(slide2, 0.5, 0.3, 6, 0.5, f"案例 {case['num']} ─ 操作指令", 16, CYAN, True)
    add_text_box(slide2, 0.5, 0.7, 12, 0.8, case["title"], 28, WHITE, True)

    # Prompt 區塊
    add_rounded_rect(slide2, 0.5, 1.6, 12.3, 5.2, CARD_BG)
    tf3 = add_text_box(slide2, 0.8, 1.7, 11.8, 0.4, "📝 複製以下指令，貼到 AI 工具中：", 16, ORANGE, True)

    prompt_lines = case["prompt"].split("\n")
    for line in prompt_lines:
        add_paragraph(tf3, line, 14, WHITE, False, 2)

    add_text_box(slide2, 0.5, 6.9, 12, 0.4, "💡 Prompt 可依實際物件修改關鍵字，不需要逐字照抄", 13, GRAY)

# ========== SOP 流程圖頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "AI 簡報製作 SOP", 32, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.2, 12, 0.5, "從接到客戶到成交後，每個環節都有對應的簡報", 18, GRAY, False, PP_ALIGN.CENTER)

sop_lines = [
    "接到新客戶 → 案例 02 客戶提案簡報",
    "    ↓",
    "帶看前 → 案例 01 社區介紹簡報",
    "    ↓",
    "每月定期 → 案例 03 區域月報簡報",
    "    ↓",
    "投資客需求 → 案例 05 投資分析 / 案例 07 店面提案",
    "    ↓",
    "成交後 → 案例 09 成交報告 → 分享給客戶的朋友（轉介紹）",
    "    ↓",
    "開發屋主 → 案例 04 委售提案簡報",
    "    ↓",
    "店長管理 → 案例 10 團隊季報 → 每季提交",
]
tf_sop = add_text_box(slide, 1, 2.0, 11, 4.5, "", 15, WHITE)
for line in sop_lines:
    c = CYAN if "↓" in line else (ORANGE if "店長管理" in line else WHITE)
    add_paragraph(tf_sop, line, 15, c, "店長管理" in line, 3)

add_text_box(slide, 0.5, 6.5, 12, 0.5, "💡 簡報不是給你看的，是給客戶看的。AI 做內容，你做最後修飾和故事", 16, ORANGE, True, PP_ALIGN.CENTER)

# ========== 結尾頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 2, 11, 1, "模組 2 完成！", 44, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "你已經學會用 AI 做出專業簡報", 28, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1.5, "記住：簡報不是給你看的，是給客戶看的\nAI 做內容，你做簡報的最後修飾和故事\n\n💡 進階挑戰：把月報簡報+月報數據串成「每月自動產出」", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6, 11, 0.5, "講師：阿峰老師 ｜ www.autolab.cloud", 16, CYAN, False, PP_ALIGN.CENTER)

output = "/Users/huangjingfeng/Desktop/專案/02-培訓(Training)/02-公開班(Public-Class)/03-ai-realtor-course(房仲AI課程)/cases/模組2-AI簡報製作-10案例.pptx"
prs.save(output)
print(f"完成：{output}")
print(f"共 {len(prs.slides)} 張投影片")
