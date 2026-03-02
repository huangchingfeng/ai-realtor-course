"""模組 1：AI 行情分析師 - 10 個實戰案例 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 品牌色
NAVY = RGBColor(0x0A, 0x16, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CARD_BG = RGBColor(0x0F, 0x1D, 0x32)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 0.8, 11, 1, "模組 1", 28, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "📊 AI 行情分析師", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 1, "10 個實戰案例 ─ 讓 AI 幫你做完行情功課", 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, "AI 超級房仲實戰班 ｜ 講師：阿峰老師（黃敬峰）", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "www.autolab.cloud", 14, CYAN, False, PP_ALIGN.CENTER)

# ========== 案例資料 ==========
cases = [
    {
        "num": "01",
        "title": "信義區社區行情比較",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：區域月報自動生成的數據基礎",
        "scenario": "客戶李美玲想在信義區買 3 房，問你翠亨村、信義之星、國泰信義三個社區哪個 CP 值最高。你需要用數據說服她，而不是憑感覺推薦。",
        "task": "用 AI 分析三個社區近 3 年實價登錄，比較單價趨勢、漲幅、公設比，產出完整的比較分析表",
        "tool": "Gemini（深度研究模式）",
        "upload_format": "📄 CSV 檔案",
        "upload_file": "實價登錄-信義區近3年.csv",
        "upload_note": "到內政部實價登錄網站下載 CSV，篩選信義區近 3 年資料",
        "prompt": "我是房仲，客戶想在信義區買 3 房，正在比較翠亨村、信義之星、國泰信義這三個社區。\n\n以下是信義區近 3 年的實價登錄資料（CSV 檔）。\n\n請幫我分析：\n1. 三個社區近 3 年的成交單價趨勢（每季平均單價）\n2. 漲幅比較：哪個社區漲最多？漲幅分別是多少 %？\n3. 公設比比較：實際可用坪數差異\n4. 樓層價差：低樓層 vs 高樓層的價差幅度\n5. 車位行情：平面車位 vs 機械車位的價格\n6. CP 值排名：綜合以上數據，哪個社區 CP 值最高？\n\n請用表格呈現比較結果，最後給出明確的推薦建議和理由。",
    },
    {
        "num": "02",
        "title": "板橋新板特區 2 房行情",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：結合模組4帶看後分析使用",
        "scenario": "首購族王大明預算 1,500 萬，想在板橋新板特區買 2 房。他在 591 上看了幾間，但不知道價格合不合理，怕買貴。",
        "task": "上傳板橋實價登錄和 591 物件截圖，讓 AI 分析 2 房成交行情和合理出價區間",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案 + 🖼️ 591 物件截圖",
        "upload_file": "實價登錄-板橋區近2年.csv + 591物件截圖",
        "upload_note": "下載板橋實價登錄 CSV，搭配 591 上有興趣的物件截圖一起上傳",
        "prompt": "我是房仲，客戶是首購族，預算 1,500 萬，想在板橋新板特區買 2 房。\n\n以下是板橋區近 2 年的實價登錄資料（CSV 檔）和他在 591 看到的物件截圖。\n\n請幫我：\n1. 篩選新板特區 2 房（20-30 坪）的成交紀錄\n2. 分析成交單價區間（最低、最高、平均）\n3. 1,500 萬預算在新板特區能買到什麼條件的 2 房？\n4. 591 截圖上這間物件的開價合理嗎？（跟實價登錄比較）\n5. 建議出價區間是多少？議價空間大概幾 %？\n\n請用表格呈現分析結果，並給出具體的議價建議。",
    },
    {
        "num": "03",
        "title": "中山區套房投資報酬率分析",
        "level": "🟡 中級",
        "system_note": "🔗 系統化：區域月報的投資版模組",
        "scenario": "投資客張先生想在中山區買套房收租，問你哪個社區投報率最高。他手上有 800 萬現金，想找穩定收租的標的。",
        "task": "分析中山區套房的成交價與租金行情，計算各社區的投資報酬率",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案（實價登錄 + 591 租金行情）",
        "upload_file": "實價登錄-中山區套房.csv + 591租金行情-中山區.csv",
        "upload_note": "需要兩份資料：實價登錄的買賣資料 + 591 上的租金行情",
        "prompt": "我是房仲，客戶是投資客，預算 800 萬，想在中山區買套房收租。\n\n以下是中山區套房的實價登錄（買賣）和 591 租金行情資料。\n\n請幫我分析：\n1. 中山區套房（10-15 坪）近 1 年的成交價區間\n2. 同區域套房的月租金行情（依地段分區）\n3. 計算各區段的毛投資報酬率（年租金 ÷ 買入價）\n4. 扣除管理費、稅費、空租期後的淨投資報酬率\n5. 800 萬預算推薦買哪一區段？為什麼？\n6. 投資風險提醒（房屋稅、租客風險、未來轉手性）\n\n請用表格呈現各區段的投報率比較，最後給出明確的投資建議。",
    },
    {
        "num": "04",
        "title": "大安區 vs 內湖區 比較分析",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "客戶陳太太一家四口想換屋，預算 3,000 萬。她老公想住大安區（上班方便），她想住內湖（有大公園、生活機能好）。兩人僵持不下，問你怎麼選。",
        "task": "用 AI 做雙區全面比較分析，幫客戶用數據做決策",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個案例不需要準備資料，直接跟 AI 對話即可",
        "prompt": "我是房仲，客戶一家四口想換屋，預算 3,000 萬。老公想住大安區，太太想住內湖區。\n\n請幫我做兩個區域的全面比較分析：\n\n1. 房價比較：3,000 萬在兩區分別能買到什麼條件的房子？（坪數、屋齡、格局）\n2. 生活機能：超市、市場、餐廳、公園的密度和品質\n3. 交通便利性：捷運站數量、到信義區/內科的通勤時間\n4. 學區比較：國小、國中的學區排名和升學率\n5. 增值潛力：近 5 年漲幅、未來重大建設（捷運延伸、都更）\n6. 居住品質：空氣品質、噪音、綠地面積\n\n請用比較表格呈現，最後根據「一家四口、有小孩」的情境給出推薦建議。\n如果可以，建議一個「折衷方案」（兩區優點都能兼顧的區域）。",
    },
    {
        "num": "05",
        "title": "林口重劃區未來發展潛力",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "客戶小夫妻預算有限，想在林口買預售屋。但他們擔心林口是不是真的會發展起來，怕買了以後房價跌。",
        "task": "用 Perplexity 即時搜尋林口最新建設進度，分析未來發展潛力",
        "tool": "Perplexity",
        "upload_format": "💬 不需要上傳，用 Perplexity 即時搜尋",
        "upload_file": "不需要檔案 — 直接在 Perplexity 輸入以下 Prompt",
        "upload_note": "Perplexity 會即時搜尋最新資訊，適合查詢建設進度和新聞",
        "prompt": "我是房仲，客戶想在林口重劃區買預售屋，請幫我搜尋並分析：\n\n1. 林口重劃區最新的重大建設進度：\n   - 機場捷運 A7/A8/A9 站周邊開發狀況\n   - 三井 Outlet 擴建計畫\n   - 國際媒體園區進度\n   - 長庚醫療園區擴展\n   - 新增的學校、公園、商場\n\n2. 交通發展：\n   - 機場捷運班次和搭乘人數趨勢\n   - 國道一號/二號交流道改善工程\n   - 未來是否有新的交通建設\n\n3. 房價走勢：\n   - 林口近 3 年的房價趨勢（跟新北市平均比較）\n   - 預售屋 vs 中古屋的價差\n   - 跟周邊區域（桃園龜山、泰山）的價差\n\n4. 投資建議：\n   - 現在進場買預售屋的優缺點\n   - 哪一區段最有潛力？\n   - 風險提醒（供給量、空屋率）\n\n請附上資料來源連結。",
    },
    {
        "num": "06",
        "title": "竹北高鐵特區 vs 關埔重劃區",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "科技業客戶在竹科上班，想在竹北買房。高鐵特區和關埔重劃區都有看，兩邊各有優缺點，猶豫不決。",
        "task": "雙工具協作：用 Perplexity 查最新開發資訊，再用 Gemini 深度分析兩區差異",
        "tool": "Gemini + Perplexity",
        "upload_format": "💬 不需要上傳，雙工具搭配使用",
        "upload_file": "不需要檔案 — 先用 Perplexity 搜尋，再貼到 Gemini 分析",
        "upload_note": "Step 1：用 Perplexity 搜尋兩區最新資訊　Step 2：把結果貼到 Gemini 深度分析",
        "prompt": "【Step 1 - Perplexity 搜尋】\n搜尋竹北高鐵特區和關埔重劃區的最新開發資訊，包括：\n- 最新建案和成交行情\n- 重大建設進度（台大醫院竹北分院、AI 智慧園區等）\n- 學區和生活機能最新變化\n\n【Step 2 - Gemini 深度分析】\n我是房仲，客戶在竹科上班，想在竹北買 3 房，預算 2,500 萬。\n以下是我搜尋到的兩區最新資訊：（貼上 Perplexity 結果）\n\n請深度分析：\n1. 房價比較：2,500 萬在兩區能買到什麼？\n2. 通勤分析：到竹科的開車/公車/自行車通勤時間\n3. 生活機能：超市、餐廳、公園、醫療的成熟度\n4. 學區品質：國小、國中的評價\n5. 增值潛力：哪區未來 5 年漲幅更大？為什麼？\n6. 適合族群：什麼樣的人適合住高鐵特區？什麼人適合關埔？\n\n請用比較表格呈現，最後給出推薦。",
    },
    {
        "num": "07",
        "title": "商辦租賃行情分析",
        "level": "🟡 中級",
        "system_note": "",
        "scenario": "企業客戶要在信義計畫區租 50 坪辦公室，預算每月 15 萬以內。他問你哪棟大樓最划算，需要數據比較。",
        "task": "分析信義計畫區商辦租金行情，找出最佳性價比的選擇",
        "tool": "ChatGPT",
        "upload_format": "📄 CSV 檔案",
        "upload_file": "591商辦租金-信義計畫區.csv",
        "upload_note": "到 591 商辦租屋頁面，篩選信義計畫區資料後匯出",
        "prompt": "我是房仲，企業客戶要在信義計畫區租辦公室，需求：50 坪左右，月租預算 15 萬以內。\n\n以下是信義計畫區的商辦租金資料（CSV 檔）。\n\n請幫我分析：\n1. 各棟商辦大樓的平均月租金（每坪）\n2. 坪效分析：實際使用面積 vs 權狀面積的比例\n3. 管理費比較：含管理費後的實際每坪成本\n4. 含稅價計算：加上營業稅後的每月實際支出\n5. 推薦 TOP 3：15 萬預算內最划算的選擇\n\n每個推薦要說明：\n- 大樓名稱、地址、樓層\n- 每月租金（含管理費、含稅）\n- 優點和缺點\n\n請用表格呈現比較結果。",
    },
    {
        "num": "08",
        "title": "法拍屋風險評估",
        "level": "🔴 進階",
        "system_note": "",
        "scenario": "投資客老陳看到一間法拍屋底價只有市價的 7 折，心動想標。但他知道法拍屋水很深，問你幫他評估風險。",
        "task": "分析法拍屋的投標風險和預期報酬，幫客戶做出投資決策",
        "tool": "ChatGPT",
        "upload_format": "🖼️ 法拍屋公告截圖或 PDF",
        "upload_file": "法拍公告-XX地方法院-案號XXX.pdf",
        "upload_note": "到司法院法拍屋查詢系統下載公告，截圖或存 PDF 上傳",
        "prompt": "我是房仲，客戶想標一間法拍屋，以下是法拍公告資料。\n\n請幫我做完整的風險評估：\n\n1. 基本資訊整理：\n   - 底價、拍次、坪數、屋齡、格局\n   - 底價相當於市價幾折？\n\n2. 風險分析（請逐項評估）：\n   - 產權風險：是否有其他抵押權、地上權、租約？\n   - 佔用風險：是否有人佔用？點交 or 不點交？\n   - 瑕疵風險：是否有漏水、結構問題的跡象？\n   - 稅費風險：增值稅、契稅、代標費等額外成本\n   - 貸款風險：法拍屋貸款成數通常較低\n\n3. 投資報酬計算：\n   - 預估總成本（得標價 + 稅費 + 裝修 + 代標費）\n   - 預估市場售價（整理後轉賣）\n   - 預估投資報酬率\n\n4. 投標建議：\n   - 建議投標價格區間\n   - 風險等級（低/中/高）\n   - 最終建議：標 or 不標？為什麼？",
    },
    {
        "num": "09",
        "title": "租買比分析（首購族決策）",
        "level": "🟢 新手",
        "system_note": "",
        "scenario": "客戶小陳 28 歲，月薪 5 萬，目前在台北租房月租 1.5 萬。他一直在猶豫到底要繼續租還是咬牙買房，問你怎麼算比較划算。",
        "task": "用數據幫客戶做租買 30 年的完整比較分析",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "這個案例不需要準備資料，直接跟 AI 對話即可",
        "prompt": "我是房仲，客戶 28 歲，月薪 5 萬，目前在台北租房月租 1.5 萬。\n他在考慮是否要買房，請幫我做 30 年的租買比較分析。\n\n買房方案假設：\n- 買 1,000 萬的房子（新北市 2 房）\n- 自備款 200 萬（頭期款 + 規費）\n- 貸款 800 萬，30 年，利率 2.1%\n\n請計算並比較：\n\n1. 每月支出比較：\n   - 租房：月租金（假設每年漲 2%）\n   - 買房：月付房貸 + 管理費 + 房屋稅 + 維修費\n\n2. 30 年總支出比較：\n   - 租房 30 年總花費\n   - 買房 30 年總花費（含利息）\n\n3. 資產累積比較：\n   - 租房 30 年後：手上有多少錢？（假設省下的錢拿去投資，年報酬 5%）\n   - 買房 30 年後：房子值多少？（假設房價年漲 2%）\n\n4. 損益平衡點：買房幾年後開始「比租房划算」？\n\n5. 決策建議：根據他的條件，租還是買？\n\n請用表格和數字呈現，讓客戶一目了然。",
    },
    {
        "num": "10",
        "title": "房貸利率敏感度分析",
        "level": "🟢 新手",
        "system_note": "🔗 系統化：結合月報，利率變動時自動推播分析給客戶",
        "scenario": "央行剛宣布升息半碼，你的客戶群組炸鍋了，好幾個人問你「利率升了對我影響多大？」你需要快速產出一份分析讓大家安心。",
        "task": "試算不同貸款金額在各利率水準下的月付金差異，做成一目了然的比較表",
        "tool": "ChatGPT",
        "upload_format": "💬 直接打字輸入（不需上傳檔案）",
        "upload_file": "不需要檔案 — 直接在 ChatGPT 輸入以下 Prompt",
        "upload_note": "每次央行升降息時都可以用這個 Prompt，秒出分析",
        "prompt": "我是房仲，央行剛宣布升息，客戶們都在問利率影響。\n\n請幫我製作一份「房貸利率敏感度分析表」：\n\n貸款金額：分別計算 1,000 萬 / 1,500 萬 / 2,000 萬\n貸款年限：30 年\n利率範圍：1.5% / 1.75% / 2.0% / 2.1% / 2.25% / 2.5% / 2.75% / 3.0%\n\n請計算每個組合的：\n1. 每月應繳金額\n2. 30 年總利息\n3. 每升息 0.25%，月付金增加多少？\n4. 每升息 0.25%，30 年多付多少利息？\n\n最後請幫我寫一段「客戶安心話術」：\n- 用白話文解釋升息的實際影響（不要用專業術語）\n- 舉例：「以 1,500 萬貸款來說，升息半碼每月只多付 XX 元，大概就是每天少喝一杯咖啡」\n- 語氣要讓客戶感覺「影響沒有想像中大」，不用太恐慌\n\n表格要清楚好讀，我要直接截圖傳給客戶。",
    },
]

# ========== 生成案例投影片 ==========
for case in cases:
    # 情境頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # 頂部標籤（含 level）
    add_text_box(slide, 0.5, 0.3, 5, 0.5, f"案例 {case['num']}  {case['level']}", 16, CYAN, True)
    add_text_box(slide, 0.5, 0.7, 12, 0.8, case["title"], 36, WHITE, True)

    # 左半邊：情境 + 任務
    tf = add_text_box(slide, 0.5, 1.8, 6, 0.5, "📋 情境", 20, ORANGE, True)
    add_paragraph(tf, case["scenario"], 17, WHITE, False, 8)
    add_paragraph(tf, "", 10, WHITE, False, 12)
    add_paragraph(tf, "🎯 任務目標", 20, ORANGE, True, 8)
    add_paragraph(tf, case["task"], 17, WHITE, False, 8)

    # 右半邊：工具 + 上傳資料
    tf2 = add_text_box(slide, 7, 1.8, 6, 0.5, "🔧 使用工具", 20, ORANGE, True)
    add_paragraph(tf2, case["tool"], 17, CYAN, True, 8)
    add_paragraph(tf2, "", 10, WHITE, False, 12)
    add_paragraph(tf2, "📎 要準備的資料", 20, ORANGE, True, 8)
    add_paragraph(tf2, f"格式：{case['upload_format']}", 15, GRAY, False, 8)
    add_paragraph(tf2, f"檔案：{case['upload_file']}", 15, WHITE, False, 4)
    add_paragraph(tf2, f"💡 {case['upload_note']}", 13, GRAY, False, 8)
    if case.get('system_note'):
        add_paragraph(tf2, case['system_note'], 13, CYAN, False, 10)

    # 操作頁（Prompt）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY)

    add_text_box(slide2, 0.5, 0.3, 6, 0.5, f"案例 {case['num']} ─ 操作指令", 16, CYAN, True)
    add_text_box(slide2, 0.5, 0.7, 12, 0.8, case["title"], 28, WHITE, True)

    # Prompt 區塊
    add_rounded_rect(slide2, 0.5, 1.6, 12.3, 5.2, CARD_BG)
    tf3 = add_text_box(slide2, 0.8, 1.7, 11.8, 0.4, "📝 複製以下指令，貼到 AI 工具中：", 16, ORANGE, True)

    prompt_lines = case["prompt"].split("\n")
    for line in prompt_lines:
        add_paragraph(tf3, line, 14, WHITE, False, 2)

    add_text_box(slide2, 0.5, 6.9, 12, 0.4, "💡 Prompt 可依實際物件修改關鍵字，不需要逐字照抄", 13, GRAY)

# ========== SOP 流程圖頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 0.5, 0.3, 12, 0.8, "AI 行情分析 SOP", 32, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.2, 12, 0.5, "從數據收集到客戶決策，每一步都有 AI 幫你", 18, GRAY, False, PP_ALIGN.CENTER)

sop_lines = [
    "每月 1 日：下載實價登錄 CSV → 案例 01 社區行情比較",
    "    ↓",
    "案例 02-04 分析不同客戶的購屋需求",
    "    ↓",
    "案例 05-06 查詢重劃區和新興區域行情",
    "    ↓",
    "案例 09 租買比分析 → 幫客戶做決策",
    "    ↓",
    "案例 10 利率敏感度 → 利率變動時即時推播",
    "    ↓",
    "所有分析結果 → 模組 2 做成簡報 → 模組 3 拆成社群內容",
]
tf_sop = add_text_box(slide, 1, 2.0, 11, 4.5, "", 15, WHITE)
for line in sop_lines:
    c = CYAN if "↓" in line else (ORANGE if "每月" in line else WHITE)
    add_paragraph(tf_sop, line, 15, c, "每月" in line, 3)

add_text_box(slide, 0.5, 6.5, 12, 0.5, "💡 數據是房仲最強的武器。AI 幫你分析，你負責解讀和說故事", 16, ORANGE, True, PP_ALIGN.CENTER)

# ========== 結尾頁 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, 1, 2, 11, 1, "模組 1 完成！", 44, CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "你已經學會用 AI 幫客戶做行情功課", 28, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 1.5, "記住：數據是房仲最強的武器\nAI 幫你分析，你負責解讀和說故事\n\n💡 進階挑戰：把案例 01+03+10 串成「區域月報自動生成」", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6, 11, 0.5, "講師：阿峰老師 ｜ www.autolab.cloud", 16, CYAN, False, PP_ALIGN.CENTER)

output = "/Users/huangjingfeng/Desktop/專案/02-培訓(Training)/02-公開班(Public-Class)/03-ai-realtor-course(房仲AI課程)/cases/模組1-AI行情分析師-10案例.pptx"
prs.save(output)
print(f"完成：{output}")
print(f"共 {len(prs.slides)} 張投影片")
